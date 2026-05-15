from __future__ import annotations

import tempfile
import unittest
import warnings
from io import StringIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

import cv2
import fitz
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Pt

import pdf2ppt.inpainting_engines as inpainting_engines
import pdf2ppt.inpainting_masks as inpainting_masks
from pdf2ppt.cli import build_parser, build_progress_callback, format_progress_line, main
from pdf2ppt.inpainting_overlay import (
    _apply_targeted_footer_label_color_correction,
    _apply_targeted_file_back_color_correction,
    _neutralize_protected_table_lines,
    _restore_protected_table_lines,
    resolve_background_inpainting_engine,
)
from pdf2ppt.inpainting_masks import refine_text_mask_for_inpainting
from pdf2ppt.models import QualityScore, TextBlock
from pdf2ppt.ocr import build_local_ocr_model_kwargs, ensure_local_model_dir, suppress_known_paddle_runtime_warnings
from pdf2ppt.pipeline import (
    BackgroundInpaintingError,
    analyze_page,
    build_text_fit_debug_entry,
    build_mask_shapes,
    build_text_mask_image,
    classify_text_script,
    ConversionOptions,
    default_font_family,
    estimate_font_size,
    estimate_background_complexity,
    estimate_text_bold,
    estimate_text_color,
    fit_text_frame,
    filter_suspicious_ocr_blocks,
    measure_text_dimensions,
    OpenCvFastInpaintingEngine,
    pil_to_image_bytes,
    promote_ocr_bold_blocks,
    resolve_background_render_dpi,
    resolve_ocr_fit_max_size,
    resolve_vertical_anchor,
    PageSignals,
    choose_background_mode,
    classify_page,
    intersection_ratio,
    mask_text_regions_with_white_boxes,
    render_overlay_background,
    select_text_blocks,
    should_wrap_text_block,
    enrich_ocr_blocks,
)
from pdf2ppt.ppt_render import add_text_block, resolve_fallback_font_size_pt
from pdf2ppt.text_style import estimate_text_style


class ClassificationTests(unittest.TestCase):
    def test_classify_digital_page(self) -> None:
        signals = PageSignals(
            native_char_count=320,
            native_text_area_ratio=0.08,
            image_area_ratio=0.15,
            drawing_count=3,
        )
        self.assertEqual(classify_page(signals), "digital")

    def test_classify_scanned_page(self) -> None:
        signals = PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=0.92,
            drawing_count=0,
        )
        self.assertEqual(classify_page(signals), "scanned")

    def test_classify_hybrid_page(self) -> None:
        signals = PageSignals(
            native_char_count=20,
            native_text_area_ratio=0.02,
            image_area_ratio=0.85,
            drawing_count=1,
        )
        self.assertEqual(classify_page(signals), "hybrid")


class BlockSelectionTests(unittest.TestCase):
    def test_hybrid_selection_filters_duplicate_ocr(self) -> None:
        native = [
            TextBlock(
                id="n1",
                source="native",
                bbox=(10, 10, 100, 30),
                text="Hello",
                confidence=0.99,
            )
        ]
        ocr = [
            TextBlock(
                id="o1",
                source="ocr",
                bbox=(12, 12, 98, 28),
                text="Hello",
                confidence=0.9,
            ),
            TextBlock(
                id="o2",
                source="ocr",
                bbox=(120, 40, 200, 80),
                text="World",
                confidence=0.9,
            ),
        ]
        selected = select_text_blocks("hybrid", native, ocr)
        self.assertEqual([block.id for block in selected], ["n1", "o2"])

    def test_intersection_ratio(self) -> None:
        ratio = intersection_ratio((0, 0, 10, 10), (5, 5, 15, 15))
        self.assertAlmostEqual(ratio, 0.25)

    def test_filter_suspicious_ocr_blocks_drops_large_low_confidence_short_text(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_bad",
                source="ocr",
                bbox=(600, 350, 780, 590),
                text="具",
                confidence=0.15,
            ),
            TextBlock(
                id="ocr_good",
                source="ocr",
                bbox=(200, 360, 520, 400),
                text="縮小至50-100份候選文檔",
                confidence=0.94,
            ),
        ]
        filtered = filter_suspicious_ocr_blocks(blocks, fitz.Rect(0, 0, 1376, 768))
        self.assertEqual([block.id for block in filtered], ["ocr_good"])

    def test_filter_suspicious_ocr_blocks_keeps_small_low_confidence_short_text(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_small",
                source="ocr",
                bbox=(100, 100, 130, 124),
                text="A",
                confidence=0.2,
            )
        ]
        filtered = filter_suspicious_ocr_blocks(blocks, fitz.Rect(0, 0, 1376, 768))
        self.assertEqual([block.id for block in filtered], ["ocr_small"])


class BackgroundModeTests(unittest.TestCase):
    def test_choose_elements_for_clean_digital_page(self) -> None:
        quality = QualityScore(
            text_confidence=0.98,
            layout_overlap_score=0.95,
            style_recovery_score=1.0,
            editable_ratio=1.0,
        )
        mode, reason = choose_background_mode(
            page_kind="digital",
            quality=quality,
            has_text=True,
            has_visuals=False,
        )
        self.assertEqual(mode, "elements")
        self.assertIsNone(reason)

    def test_choose_full_page_for_low_editable_score(self) -> None:
        quality = QualityScore(
            text_confidence=0.4,
            layout_overlap_score=0.3,
            style_recovery_score=0.2,
            editable_ratio=0.2,
        )
        mode, reason = choose_background_mode(
            page_kind="scanned",
            quality=quality,
            has_text=False,
            has_visuals=True,
        )
        self.assertEqual(mode, "full-page")
        self.assertIsNotNone(reason)

    def test_mask_text_regions_with_white_boxes(self) -> None:
        image = Image.new("RGB", (100, 60), color=(20, 30, 40))
        blocks = [
            TextBlock(
                id="ocr_1",
                source="ocr",
                bbox=(10, 10, 30, 25),
                text="demo",
                confidence=0.9,
                image_bbox=(10, 10, 30, 25),
            )
        ]
        masked = mask_text_regions_with_white_boxes(image, blocks, fitz.Rect(0, 0, 100, 60))
        self.assertEqual(masked.getpixel((0, 0)), (20, 30, 40))
        self.assertEqual(masked.getpixel((20, 16)), (255, 255, 255))

    def test_mask_prefers_image_polygon(self) -> None:
        image = Image.new("RGB", (60, 60), color=(10, 10, 10))
        blocks = [
            TextBlock(
                id="ocr_2",
                source="ocr",
                bbox=(0, 0, 60, 60),
                text="demo",
                confidence=0.9,
                image_polygon=((10, 10), (40, 10), (40, 20), (10, 20)),
            )
        ]
        masked = mask_text_regions_with_white_boxes(image, blocks, fitz.Rect(0, 0, 60, 60))
        self.assertEqual(masked.getpixel((15, 15)), (255, 255, 255))
        self.assertEqual(masked.getpixel((45, 25)), (10, 10, 10))

    def test_build_mask_shapes_uses_exact_polygon(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_3",
                source="ocr",
                bbox=(0, 0, 60, 60),
                text="demo",
                confidence=0.9,
                image_polygon=((10, 10), (40, 10), (40, 20), (10, 20)),
            )
        ]
        shapes = build_mask_shapes(blocks, (60, 60), fitz.Rect(0, 0, 60, 60))
        self.assertEqual(shapes, [{"kind": "polygon", "points": [(10, 10), (40, 10), (40, 20), (10, 20)]}])

    def test_build_mask_shapes_scales_page_space_polygon_to_target_image(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_3b",
                source="ocr",
                bbox=(0, 0, 60, 60),
                text="demo",
                confidence=0.9,
                image_polygon=((20, 20), (80, 20), (80, 40), (20, 40)),
            )
        ]
        shapes = build_mask_shapes(blocks, (50, 50), fitz.Rect(0, 0, 100, 100))
        self.assertEqual(shapes, [{"kind": "polygon", "points": [(10, 10), (40, 10), (40, 20), (10, 20)]}])

    def test_build_text_mask_image_applies_padding(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_4",
                source="ocr",
                bbox=(10, 10, 20, 20),
                text="demo",
                confidence=0.9,
                image_bbox=(10, 10, 20, 20),
            )
        ]
        mask = build_text_mask_image(blocks, (40, 40), fitz.Rect(0, 0, 40, 40), padding_px=2)
        self.assertEqual(mask.getpixel((10, 10)), 255)
        self.assertEqual(mask.getpixel((8, 10)), 255)
        self.assertEqual(mask.getpixel((5, 10)), 0)

    def test_refine_text_mask_for_inpainting_preserves_table_line_in_padded_fringe(self) -> None:
        image = Image.new("RGB", (80, 80), color=(255, 255, 255))
        draw = ImageDraw.Draw(image)
        draw.line((0, 18, 79, 18), fill=(90, 90, 90), width=2)
        blocks = [
            TextBlock(
                id="ocr_refine_1",
                source="ocr",
                bbox=(20, 20, 60, 40),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 60, 40),
            )
        ]

        mask = refine_text_mask_for_inpainting(image, blocks, (80, 80), fitz.Rect(0, 0, 80, 80), padding_px=4)

        self.assertEqual(mask.getpixel((30, 18)), 0)
        self.assertEqual(mask.getpixel((30, 22)), 255)

    def test_refine_text_mask_for_inpainting_keeps_padding_when_no_table_line_exists(self) -> None:
        image = Image.new("RGB", (80, 80), color=(255, 255, 255))
        blocks = [
            TextBlock(
                id="ocr_refine_2",
                source="ocr",
                bbox=(20, 20, 60, 40),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 60, 40),
            )
        ]

        mask = refine_text_mask_for_inpainting(image, blocks, (80, 80), fitz.Rect(0, 0, 80, 80), padding_px=4)

        self.assertEqual(mask.getpixel((18, 30)), 255)

    def test_refine_text_mask_for_inpainting_preserves_table_line_inside_ocr_boundary_band(self) -> None:
        image = Image.new("RGB", (80, 80), color=(255, 255, 255))
        draw = ImageDraw.Draw(image)
        draw.line((20, 20, 60, 20), fill=(90, 90, 90), width=2)
        blocks = [
            TextBlock(
                id="ocr_refine_3",
                source="ocr",
                bbox=(20, 20, 60, 40),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 60, 40),
            )
        ]

        mask = refine_text_mask_for_inpainting(image, blocks, (80, 80), fitz.Rect(0, 0, 80, 80), padding_px=0)

        self.assertEqual(mask.getpixel((30, 20)), 0)
        self.assertEqual(mask.getpixel((30, 30)), 255)

    def test_render_overlay_background_auto_uses_opencv_fast_for_small_mask(self) -> None:
        image = Image.new("RGB", (60, 40), color=(35, 45, 55))
        blocks = [
            TextBlock(
                id="ocr_5",
                source="ocr",
                bbox=(20, 12, 35, 24),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 12, 35, 24),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="auto",
            inpaint_padding_px=0,
            inpaint_max_area_ratio=0.2,
        )
        result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 60, 40), options=options)
        self.assertEqual(result.engine_name, "opencv-fast")
        self.assertIn("opencv-fast", result.note or "")
        pixel = result.image.getpixel((25, 18))
        self.assertTrue(all(abs(channel - expected) <= 4 for channel, expected in zip(pixel, (35, 45, 55))))

    def test_resolve_background_inpainting_engine_returns_lama_for_explicit_request(self) -> None:
        image = Image.new("RGB", (60, 40), color=(35, 45, 55))
        mask_image = Image.new("L", (60, 40), color=0)
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="lama-onnx-cuda",
            inpaint_model_root=Path("model/lama"),
            inpaint_onnx_cuda_provider="CUDAExecutionProvider",
            inpaint_onnx_execution_mode="parallel",
            inpaint_max_side_px=1024,
        )

        engine, note = resolve_background_inpainting_engine(image, mask_image, options)

        self.assertIsInstance(engine, inpainting_engines.LamaOnnxCudaInpaintingEngine)
        self.assertIn("lama-onnx-cuda", note)
        self.assertEqual(engine.max_side_px, 1024)

    def test_render_overlay_background_explicit_lama_is_strict(self) -> None:
        image = Image.new("RGB", (60, 40), color=(35, 45, 55))
        blocks = [
            TextBlock(
                id="ocr_lama_strict",
                source="ocr",
                bbox=(20, 12, 35, 24),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 12, 35, 24),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="lama-onnx-cuda",
            inpaint_model_root=Path("model/lama"),
        )

        with patch.object(
            inpainting_engines.LamaOnnxCudaInpaintingEngine,
            "inpaint",
            side_effect=BackgroundInpaintingError("missing model"),
        ), patch.object(inpainting_engines.WhiteBoxInpaintingEngine, "inpaint") as white_box_mock:
            with self.assertRaisesRegex(BackgroundInpaintingError, "missing model"):
                render_overlay_background(image, blocks, fitz.Rect(0, 0, 60, 40), options=options)

        white_box_mock.assert_not_called()

    def test_lama_engine_requires_existing_model_root(self) -> None:
        engine = inpainting_engines.LamaOnnxCudaInpaintingEngine(model_root=Path("missing-lama-model"))

        with self.assertRaisesRegex(BackgroundInpaintingError, "does not exist"):
            engine.inpaint(Image.new("RGB", (20, 20), color=(10, 20, 30)), Image.new("L", (20, 20), color=255))

    def test_lama_engine_reports_missing_onnxruntime_gpu_dependency(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            model_path = Path(temp_dir) / "lama_fp16.onnx"
            model_path.write_bytes(b"fake-onnx")
            engine = inpainting_engines.LamaOnnxCudaInpaintingEngine(model_root=model_path)

            with patch("pdf2ppt.inpainting_engines.importlib.import_module", side_effect=ImportError("missing ort")):
                with self.assertRaisesRegex(BackgroundInpaintingError, "onnxruntime-gpu"):
                    engine.inpaint(Image.new("RGB", (20, 20), color=(10, 20, 30)), Image.new("L", (20, 20), color=255))

    def test_lama_model_input_mapping_prefers_named_inputs(self) -> None:
        fake_session = SimpleNamespace(
            get_inputs=lambda: [
                SimpleNamespace(name="image"),
                SimpleNamespace(name="mask"),
            ]
        )
        image_tensor = np.zeros((1, 3, 8, 8), dtype=np.float32)
        mask_tensor = np.ones((1, 1, 8, 8), dtype=np.float32)

        resolved = inpainting_engines._build_lama_model_inputs(fake_session, image_tensor, mask_tensor)

        self.assertIs(resolved["image"], image_tensor)
        self.assertIs(resolved["mask"], mask_tensor)

    def test_lama_fixed_input_size_detects_official_opencv_shape(self) -> None:
        fake_session = SimpleNamespace(
            get_inputs=lambda: [
                SimpleNamespace(name="image", shape=["batch", 3, 512, 512]),
                SimpleNamespace(name="mask", shape=["batch", 1, 512, 512]),
            ]
        )

        fixed_input_size = inpainting_engines._resolve_lama_fixed_input_size(fake_session)

        self.assertEqual(fixed_input_size, (512, 512))

    def test_lama_fixed_input_adapter_resizes_to_model_shape(self) -> None:
        source_rgb = np.zeros((864, 1536, 3), dtype=np.uint8)
        mask_array = np.zeros((864, 1536), dtype=np.uint8)
        mask_array[120:240, 300:520] = 255

        resized_rgb, resized_mask, adapted = inpainting_engines._fit_lama_inputs_to_model(
            source_rgb,
            mask_array,
            (512, 512),
        )

        self.assertTrue(adapted)
        self.assertEqual(resized_rgb.shape, (512, 512, 3))
        self.assertEqual(resized_mask.shape, (512, 512))
        self.assertGreater(np.count_nonzero(resized_mask), 0)

    def test_render_overlay_background_emits_mask_debug_images(self) -> None:
        image = Image.new("RGB", (80, 80), color=(255, 255, 255))
        blocks = [
            TextBlock(
                id="ocr_debug_masks",
                source="ocr",
                bbox=(20, 20, 60, 60),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 60, 60),
            )
        ]
        horizontal = np.zeros((80, 80), dtype=bool)
        vertical = np.zeros((80, 80), dtype=bool)
        horizontal[38:41, 12:68] = True
        vertical[12:68, 38:41] = True
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=4,
            inpaint_max_area_ratio=0.2,
        )

        with patch.object(inpainting_masks, "_detect_table_line_orientation_masks", return_value=(horizontal, vertical)):
            result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 80, 80), options=options)

        self.assertIn("raw_mask", result.debug_images)
        self.assertIn("refined_mask", result.debug_images)
        self.assertIn("table_line_mask", result.debug_images)
        self.assertIn("grid_line_mask", result.debug_images)
        self.assertIn("protected_line_mask", result.debug_images)
        self.assertIn("protected_prefill_ring_mask", result.debug_images)

    def test_render_overlay_background_skips_protected_line_mask_for_isolated_line(self) -> None:
        image = Image.new("RGB", (80, 80), color=(255, 255, 255))
        draw = ImageDraw.Draw(image)
        draw.line((20, 20, 60, 20), fill=(90, 90, 90), width=2)
        blocks = [
            TextBlock(
                id="ocr_debug_isolated_line",
                source="ocr",
                bbox=(20, 20, 60, 40),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 60, 40),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=4,
            inpaint_max_area_ratio=0.2,
        )

        result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 80, 80), options=options)

        self.assertIn("table_line_mask", result.debug_images)
        self.assertIn("grid_line_mask", result.debug_images)
        self.assertNotIn("protected_line_mask", result.debug_images)

    def test_render_overlay_background_skips_protected_line_mask_for_text_overlapping_false_grid(self) -> None:
        image = Image.new("RGB", (100, 100), color=(255, 255, 255))
        blocks = [
            TextBlock(
                id="ocr_false_grid_guard",
                source="ocr",
                bbox=(20, 20, 80, 80),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 20, 80, 80),
            )
        ]
        horizontal = np.zeros((100, 100), dtype=bool)
        vertical = np.zeros((100, 100), dtype=bool)
        horizontal[48:51, 22:79] = True
        vertical[22:79, 48:51] = True
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=4,
            inpaint_max_area_ratio=0.2,
        )

        with patch.object(inpainting_masks, "_detect_table_line_orientation_masks", return_value=(horizontal, vertical)):
            result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 100, 100), options=options)

        self.assertIn("table_line_mask", result.debug_images)
        self.assertIn("grid_line_mask", result.debug_images)
        self.assertNotIn("protected_line_mask", result.debug_images)

    def test_neutralize_and_restore_protected_table_lines_avoids_black_border_sampling(self) -> None:
        source = np.full((80, 120, 3), 248, dtype=np.uint8)
        source[20:22, 20:100] = 40
        mask_array = np.zeros((80, 120), dtype=np.uint8)
        mask_array[22:42, 24:96] = 255
        protected_line_mask = np.zeros((80, 120), dtype=bool)
        protected_line_mask[20:22, 20:100] = True

        neutralized = _neutralize_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            mask_array,
            protected_line_mask,
        )
        neutralized_array = np.array(neutralized, dtype=np.int16)
        self.assertGreater(neutralized_array[20:22, 30:90].mean(), 220.0)

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            neutralized,
            protected_line_mask,
        )
        restored_array = np.array(restored, dtype=np.int16)
        source_array = source.astype(np.int16)
        self.assertLess(np.abs(restored_array[20:22, 30:90] - source_array[20:22, 30:90]).mean(), 1.0)

    def test_restore_protected_table_lines_recovers_dark_line_fringe(self) -> None:
        source = np.full((40, 80, 3), 245, dtype=np.uint8)
        source[18:20, 10:70] = 40
        source[17, 10:70] = 78
        source[20, 10:70] = 78
        source[16, 10:70] = 96
        source[21, 10:70] = 96
        source[22, 10:70] = 110
        rendered = np.full((40, 80, 3), 245, dtype=np.uint8)
        protected_line_mask = np.zeros((40, 80), dtype=bool)
        protected_line_mask[18:20, 10:70] = True

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            Image.fromarray(rendered, mode="RGB"),
            protected_line_mask,
        )

        restored_array = np.array(restored, dtype=np.int16)
        source_array = source.astype(np.int16)
        self.assertLess(np.abs(restored_array[17, 20:60] - source_array[17, 20:60]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[20, 20:60] - source_array[20, 20:60]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[16, 20:60] - source_array[16, 20:60]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[21, 20:60] - source_array[21, 20:60]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[22, 20:60] - source_array[22, 20:60]).mean(), 1.0)
        self.assertGreater(restored_array[10:14, 20:60].mean(), 240.0)

    def test_restore_protected_table_lines_extends_low_contrast_footer_border_from_seed(self) -> None:
        width, height = 220, 140
        source = np.full((height, width, 3), 240, dtype=np.uint8)
        rendered = np.full((height, width, 3), 240, dtype=np.uint8)
        source[132:136, 30:190] = 96
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[128:130, 104:114] = True
        source[128:130, 104:114] = 72
        rendered[132:136, 30:190] = 240
        block = TextBlock(
            id="ocr_footer_1",
            source="ocr",
            bbox=(30.0, 90.0, 190.0, 130.0),
            text="高頻FAQ(98.8%)",
            confidence=0.99,
            block_role="body",
        )

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            Image.fromarray(rendered, mode="RGB"),
            protected_line_mask,
            text_blocks=[block],
            page_rect=fitz.Rect(0, 0, width, height),
        )

        restored_array = np.array(restored, dtype=np.int16)
        source_array = source.astype(np.int16)
        self.assertLess(np.abs(restored_array[132:136, 50:170] - source_array[132:136, 50:170]).mean(), 1.0)

    def test_restore_protected_table_lines_recovers_footer_light_band_without_seed(self) -> None:
        width, height = 220, 140
        source = np.full((height, width, 3), 240, dtype=np.uint8)
        rendered = np.full((height, width, 3), 240, dtype=np.uint8)
        source[83:86, 30:190] = 242
        rendered[83:86, 30:190] = 234
        source[86:89, 30:190] = 242
        rendered[86:89, 30:190] = 234
        source[120:126, 30:190] = 236
        rendered[120:126, 30:190] = 140
        source[126:130, 30:190] = 236
        source[130:134, 30:190] = 96
        rendered[126:130, 30:190] = 140
        source[131, 50:108] = 236
        rendered[131, 50:108] = 140
        rendered[130:134, 30:190] = 96
        protected_line_mask = np.zeros((height, width), dtype=bool)
        block = TextBlock(
            id="ocr_footer_2",
            source="ocr",
            bbox=(30.0, 86.0, 190.0, 126.0),
            text="同義改寫 (100%)",
            confidence=0.99,
            block_role="body",
        )

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            Image.fromarray(rendered, mode="RGB"),
            protected_line_mask,
            text_blocks=[block],
            page_rect=fitz.Rect(0, 0, width, height),
        )

        restored_array = np.array(restored, dtype=np.int16)
        source_array = source.astype(np.int16)
        self.assertLess(np.abs(restored_array[83:86, 50:170] - source_array[83:86, 50:170]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[86:89, 50:170] - source_array[86:89, 50:170]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[120:126, 50:170] - source_array[120:126, 50:170]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[126:130, 50:170] - source_array[126:130, 50:170]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[131, 50:108] - source_array[131, 50:108]).mean(), 1.0)
        self.assertLess(np.abs(restored_array[130:134, 50:170] - source_array[130:134, 50:170]).mean(), 1.0)

    def test_restore_protected_table_lines_skips_non_footer_body_blocks_for_footer_border_restore(self) -> None:
        width, height = 220, 140
        source = np.full((height, width, 3), 240, dtype=np.uint8)
        rendered = np.full((height, width, 3), 240, dtype=np.uint8)
        source[108:112, 30:190] = 96
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[100:102, 104:114] = True
        source[100:102, 104:114] = 72
        rendered[108:112, 30:190] = 240
        block = TextBlock(
            id="ocr_body_midpage_guard_1",
            source="ocr",
            bbox=(30.0, 60.0, 190.0, 100.0),
            text="Mid-page label",
            confidence=0.99,
            block_role="body",
        )

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            Image.fromarray(rendered, mode="RGB"),
            protected_line_mask,
            text_blocks=[block],
            page_rect=fitz.Rect(0, 0, width, height),
        )

        restored_array = np.array(restored, dtype=np.int16)
        self.assertGreater(restored_array[108:112, 50:170].mean(), 230.0)

    def test_restore_protected_table_lines_skips_title_blocks_for_footer_border_restore(self) -> None:
        width, height = 220, 140
        source = np.full((height, width, 3), 240, dtype=np.uint8)
        rendered = np.full((height, width, 3), 240, dtype=np.uint8)
        source[52:56, 30:190] = 96
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[30:32, 104:114] = True
        source[30:32, 104:114] = 72
        rendered[52:56, 30:190] = 240
        block = TextBlock(
            id="ocr_title_footer_guard_1",
            source="ocr",
            bbox=(30.0, 10.0, 190.0, 50.0),
            text="Main Title",
            confidence=0.99,
            block_role="title",
        )

        restored = _restore_protected_table_lines(
            Image.fromarray(source, mode="RGB"),
            Image.fromarray(rendered, mode="RGB"),
            protected_line_mask,
            text_blocks=[block],
            page_rect=fitz.Rect(0, 0, width, height),
        )

        restored_array = np.array(restored, dtype=np.int16)
        self.assertGreater(restored_array[52:56, 50:170].mean(), 230.0)

    def test_render_overlay_background_auto_falls_back_to_white_box_for_large_mask(self) -> None:
        grid_y, grid_x = np.indices((40, 60), dtype=np.uint8)
        textured = np.stack(
            [
                (grid_x * 17 + grid_y * 11) % 255,
                (grid_x * 29 + grid_y * 7) % 255,
                (grid_x * 13 + grid_y * 19) % 255,
            ],
            axis=2,
        )
        image = Image.fromarray(textured, mode="RGB")
        blocks = [
            TextBlock(
                id="ocr_6",
                source="ocr",
                bbox=(5, 5, 55, 35),
                text="demo",
                confidence=0.9,
                image_bbox=(5, 5, 55, 35),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="auto",
            inpaint_padding_px=0,
            inpaint_max_area_ratio=0.1,
        )
        result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 60, 40), options=options)
        self.assertEqual(result.engine_name, "white-box")
        self.assertIn("white-box", result.note or "")
        self.assertEqual(result.image.getpixel((20, 18)), (255, 255, 255))


class OcrStyleOptimizationTests(unittest.TestCase):
    def test_estimate_text_style_matches_individual_estimators(self) -> None:
        image = Image.new("RGB", (64, 24), color=(255, 255, 255))
        draw = ImageDraw.Draw(image)
        draw.text((6, 4), "Demo", fill=(24, 36, 48), font=ImageFont.load_default())
        gray = image.convert("L")

        combined_color, combined_bold = estimate_text_style("Demo", image, gray)

        self.assertEqual(combined_color, estimate_text_color(image, gray))
        self.assertEqual(combined_bold, estimate_text_bold("Demo", gray))

    def test_enrich_ocr_blocks_uses_combined_style_estimator_once_per_block(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_1",
                source="ocr",
                bbox=(5, 6, 45, 20),
                text="Hello",
                confidence=0.95,
            ),
            TextBlock(
                id="ocr_2",
                source="ocr",
                bbox=(10, 24, 52, 38),
                text="World",
                confidence=0.91,
            ),
        ]
        image = Image.new("RGB", (80, 50), color=(255, 255, 255))

        with (
            patch("pdf2ppt.block_analysis.assign_block_roles") as assign_mock,
            patch("pdf2ppt.block_analysis.promote_ocr_bold_blocks") as promote_mock,
            patch("pdf2ppt.block_analysis.sort_text_blocks", side_effect=lambda items: items),
            patch("pdf2ppt.block_analysis.classify_text_script", return_value="latin") as script_mock,
            patch(
                "pdf2ppt.block_analysis.estimate_text_style",
                return_value=("#102030", True),
            ) as style_mock,
            patch("pdf2ppt.block_analysis.estimate_font_size", return_value=14.0) as font_mock,
        ):
            enriched = enrich_ocr_blocks(blocks, image)

        self.assertEqual(len(enriched), 2)
        self.assertEqual(script_mock.call_count, len(blocks))
        self.assertEqual(style_mock.call_count, len(blocks))
        self.assertEqual(font_mock.call_count, len(blocks))
        self.assertEqual([block.font_color for block in enriched], ["#102030", "#102030"])
        self.assertEqual([block.bold for block in enriched], [True, True])
        self.assertEqual([block.font_size for block in enriched], [14.0, 14.0])
        for call in style_mock.call_args_list:
            self.assertEqual(call.kwargs["script"], "latin")
        for call in font_mock.call_args_list:
            self.assertEqual(call.kwargs["script"], "latin")
        assign_mock.assert_called_once()
        promote_mock.assert_called_once()

    def test_render_overlay_background_auto_uses_opencv_fast_for_large_low_texture_mask(self) -> None:
        width, height = 180, 120
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        base = np.stack(
            [
                150.0 + 0.03 * grid_x + 0.08 * grid_y,
                140.0 + 0.02 * grid_x + 0.06 * grid_y,
                130.0 + 0.01 * grid_x + 0.05 * grid_y,
            ],
            axis=2,
        ).clip(0, 255).astype("uint8")
        image = Image.fromarray(base, mode="RGB")
        blocks = [
            TextBlock(
                id="ocr_large_flat",
                source="ocr",
                bbox=(45, 25, 135, 95),
                text="demo",
                confidence=0.9,
                image_bbox=(45, 25, 135, 95),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="auto",
            inpaint_padding_px=0,
            inpaint_max_area_ratio=0.12,
        )
        result = render_overlay_background(image, blocks, fitz.Rect(0, 0, width, height), options=options)
        self.assertEqual(result.engine_name, "opencv-fast")
        self.assertIn("low-texture", result.note or "")

        result_array = np.array(result.image, dtype=np.int16)
        base_array = base.astype(np.int16)
        mask = np.zeros((height, width), dtype=bool)
        mask[25:95, 45:135] = True
        mae = np.abs(result_array - base_array)[mask].mean()
        self.assertLess(mae, 1.0)

    def test_background_complexity_detects_textured_context(self) -> None:
        image = Image.linear_gradient("L").resize((80, 60)).convert("RGB")
        mask = Image.new("L", (80, 60), 0)
        for x in range(30, 50):
            for y in range(20, 35):
                mask.putpixel((x, y), 255)
        complexity = estimate_background_complexity(image, mask)
        self.assertGreater(complexity, 0.0)

    def test_opencv_fast_restores_smooth_gradient_background(self) -> None:
        width, height = 180, 120
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        base = np.stack(
            [
                150.0 + 0.03 * grid_x + 0.08 * grid_y,
                140.0 + 0.02 * grid_x + 0.06 * grid_y,
                130.0 + 0.01 * grid_x + 0.05 * grid_y,
            ],
            axis=2,
        ).clip(0, 255).astype("uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[25:95, 45:135] = 255

        result = OpenCvFastInpaintingEngine().inpaint(
            Image.fromarray(base, mode="RGB"),
            Image.fromarray(mask, mode="L"),
        )

        result_array = np.array(result, dtype=np.int16)
        base_array = base.astype(np.int16)
        mae = np.abs(result_array - base_array)[mask > 0].mean()
        self.assertLess(mae, 1.0)

    def test_opencv_fast_skips_telea_for_smooth_gradient_background(self) -> None:
        width, height = 180, 120
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        base = np.stack(
            [
                150.0 + 0.03 * grid_x + 0.08 * grid_y,
                140.0 + 0.02 * grid_x + 0.06 * grid_y,
                130.0 + 0.01 * grid_x + 0.05 * grid_y,
            ],
            axis=2,
        ).clip(0, 255).astype("uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[25:95, 45:135] = 255

        with patch("pdf2ppt.inpainting_engines.cv2.inpaint") as inpaint_mock:
            result = OpenCvFastInpaintingEngine().inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        inpaint_mock.assert_not_called()
        result_array = np.array(result, dtype=np.int16)
        base_array = base.astype(np.int16)
        mae = np.abs(result_array - base_array)[mask > 0].mean()
        self.assertLess(mae, 1.0)

    def test_opencv_fast_uses_adaptive_telea_radius_per_component_size(self) -> None:
        width, height = 220, 160
        base = np.zeros((height, width, 3), dtype="uint8")
        for channel in range(3):
            base[:, :, channel] = np.linspace(40 + channel * 20, 180 + channel * 15, width, dtype=np.uint8)
        mask = np.zeros((height, width), dtype="uint8")
        mask[16:28, 16:28] = 255
        mask[50:122, 110:198] = 255

        captured_radii: list[float] = []

        def fake_inpaint(local_source: np.ndarray, local_mask: np.ndarray, radius: float, method: int) -> np.ndarray:
            captured_radii.append(radius)
            return local_source.copy()

        with (
            patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base.copy(), mask.copy())),
            patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=fake_inpaint),
        ):
            OpenCvFastInpaintingEngine().inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        captured_radii.sort()
        self.assertEqual(len(captured_radii), 2)
        self.assertLess(captured_radii[0], captured_radii[1])
        self.assertGreaterEqual(captured_radii[0], 1.79)
        self.assertLessEqual(captured_radii[1], 7.51)

    def test_opencv_fast_reduces_telea_radius_when_ring_has_dense_edges(self) -> None:
        width, height = 220, 160
        base = np.full((height, width, 3), 180, dtype="uint8")
        base[:, 110:, :] = 120
        for y in range(height):
            if y % 6 < 3:
                base[y, 120:180, :] = 220
        mask = np.zeros((height, width), dtype="uint8")
        mask[40:84, 24:68] = 255
        mask[40:84, 132:176] = 255

        captured_radii: list[float] = []

        def fake_inpaint(local_source: np.ndarray, local_mask: np.ndarray, radius: float, method: int) -> np.ndarray:
            captured_radii.append(radius)
            return local_source.copy()

        with (
            patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base.copy(), mask.copy())),
            patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=fake_inpaint),
        ):
            OpenCvFastInpaintingEngine(
                telea_small_component_max_span_px=20,
                telea_group_proximity_px=0,
            ).inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        self.assertEqual(len(captured_radii), 2)
        self.assertGreater(captured_radii[0], captured_radii[1])

    def test_opencv_fast_shrinks_radius_and_crop_near_protected_lines(self) -> None:
        width, height = 220, 160
        base = np.full((height, width, 3), 180, dtype="uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[40:84, 24:68] = 255
        mask[40:84, 132:176] = 255
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[32:36, 16:76] = True

        captured_calls: list[tuple[tuple[int, int], float, tuple[int, int, int, int]]] = []

        def fake_inpaint(local_source: np.ndarray, local_mask: np.ndarray, radius: float, method: int) -> np.ndarray:
            points = cv2.findNonZero(local_mask)
            bbox = cv2.boundingRect(points) if points is not None else (0, 0, 0, 0)
            captured_calls.append((local_source.shape[:2], radius, bbox))
            return local_source.copy()

        with (
            patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base.copy(), mask.copy())),
            patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=fake_inpaint),
        ):
            engine = OpenCvFastInpaintingEngine(
                telea_small_component_max_span_px=20,
                telea_group_proximity_px=0,
            )
            engine.set_protected_line_mask(protected_line_mask)
            engine.inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        self.assertEqual(len(captured_calls), 2)
        protected_call, regular_call = captured_calls
        self.assertLess(protected_call[1], regular_call[1])
        self.assertLess(protected_call[0][0] * protected_call[0][1], regular_call[0][0] * regular_call[0][1])
        protected_height, protected_width = protected_call[0]
        protected_x, protected_y, protected_bbox_width, protected_bbox_height = protected_call[2]
        top_margin = protected_y
        bottom_margin = protected_height - (protected_y + protected_bbox_height)
        self.assertLess(top_margin, bottom_margin)

    def test_inpaint_residual_components_excludes_protected_lines_from_edge_density(self) -> None:
        width, height = 220, 160
        source = np.full((height, width, 3), 180, dtype="uint8")
        residual_mask = np.zeros((height, width), dtype="uint8")
        residual_mask[40:84, 24:68] = 255
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[32:36, 16:76] = True
        source[32:36, 16:76] = 20

        with patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=lambda local_source, local_mask, radius, method: local_source.copy()):
            _, diagnostics = inpainting_engines._inpaint_residual_components(
                source,
                residual_mask,
                protected_line_mask=protected_line_mask,
                base_radius=3.0,
                min_radius=1.8,
                max_radius=7.5,
                reference_span_px=48.0,
                edge_density_threshold=0.08,
                edge_density_min_factor=0.7,
                small_component_max_span_px=20,
                group_proximity_px=0,
                group_proximity_min_scale=0.75,
                group_proximity_max_scale=2.0,
            )

        self.assertEqual(len(diagnostics), 1)
        self.assertEqual(int(diagnostics[0]["protected_nearby"]), 1)
        self.assertEqual(float(diagnostics[0]["edge_density"]), 0.0)

    def test_resolve_component_background_patch_excludes_protected_lines_from_prefill_ring(self) -> None:
        width, height = 180, 120
        source = np.full((height, width, 3), 180, dtype="uint8")
        component = np.zeros((height, width), dtype="uint8")
        component[40:80, 60:120] = 1
        protected_line_mask = np.zeros((height, width), dtype=bool)
        protected_line_mask[34:38, 48:132] = True
        source[34:38, 48:132] = 20
        gray = cv2.cvtColor(source, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 80, 160)
        kernel = np.ones((17, 17), dtype=np.uint8)
        expand_kernel = np.ones((5, 5), dtype=np.uint8)
        gap_kernel = np.ones((13, 13), dtype=np.uint8)
        captured_overlaps: list[int] = []

        def fake_fit_component_background_surface(
            local_source: np.ndarray,
            local_component: np.ndarray,
            ring_mask: np.ndarray,
            *,
            context_dilate_px: int,
            model: str,
        ) -> tuple[int, int, int, int, np.ndarray, float]:
            captured_overlaps.append(int(np.count_nonzero(ring_mask & protected_line_mask)))
            patch_height, patch_width = local_component.shape
            patch_values = np.full((patch_height, patch_width, 3), 180.0, dtype=np.float32)
            return 0, 0, patch_width, patch_height, patch_values, 0.0

        with patch(
            "pdf2ppt.inpainting_engines._fit_component_background_surface",
            side_effect=fake_fit_component_background_surface,
        ):
            inpainting_engines._resolve_component_background_patch(
                source,
                gray,
                edges,
                component,
                protected_line_mask=None,
                protected_prefill_ring_mask=None,
                kernel=kernel,
                expand_kernel=expand_kernel,
                gap_kernel=gap_kernel,
                flat_background_std_threshold=4.0,
                flat_background_edge_threshold=0.01,
                context_dilate_px=8,
                smooth_gradient_edge_threshold=0.015,
                smooth_gradient_residual_threshold=16.0,
                smooth_gradient_color_bias_max_delta=0.0,
                smooth_gradient_color_bias_residual_scale=4.0,
            )
            inpainting_engines._resolve_component_background_patch(
                source,
                gray,
                edges,
                component,
                protected_line_mask=protected_line_mask,
                protected_prefill_ring_mask=np.zeros((height, width), dtype=np.uint8),
                kernel=kernel,
                expand_kernel=expand_kernel,
                gap_kernel=gap_kernel,
                flat_background_std_threshold=4.0,
                flat_background_edge_threshold=0.01,
                context_dilate_px=8,
                smooth_gradient_edge_threshold=0.015,
                smooth_gradient_residual_threshold=16.0,
                smooth_gradient_color_bias_max_delta=0.0,
                smooth_gradient_color_bias_residual_scale=4.0,
            )

        self.assertEqual(len(captured_overlaps), 2)
        self.assertGreater(captured_overlaps[0], 0)
        self.assertEqual(captured_overlaps[1], 0)

    def test_resolve_component_background_patch_allows_relaxed_quadratic_for_wide_short_component(self) -> None:
        width, height = 420, 160
        source = np.full((height, width, 3), 180, dtype="uint8")
        component = np.zeros((height, width), dtype="uint8")
        component[70:118, 40:260] = 1
        gray = cv2.cvtColor(source, cv2.COLOR_BGR2GRAY)
        edges = np.zeros((height, width), dtype="uint8")
        kernel = np.ones((17, 17), dtype=np.uint8)
        expand_kernel = np.ones((5, 5), dtype=np.uint8)
        gap_kernel = np.ones((13, 13), dtype=np.uint8)

        expanded = cv2.dilate(component, expand_kernel, iterations=1)
        context_component = cv2.dilate(expanded, gap_kernel, iterations=1)
        ring_mask = cv2.dilate(context_component, kernel, iterations=1).astype(bool) & (~context_component.astype(bool))
        ring_points = np.argwhere(ring_mask)
        edge_count = max(1, int(round(float(np.count_nonzero(ring_mask)) * 0.02)))
        edges[tuple(ring_points[:edge_count].T)] = 255

        call_models: list[str] = []

        def fake_fit_component_background_surface(
            local_source: np.ndarray,
            local_component: np.ndarray,
            local_ring_mask: np.ndarray,
            *,
            context_dilate_px: int,
            model: str,
        ) -> tuple[int, int, int, int, np.ndarray, float] | None:
            call_models.append(model)
            patch_height, patch_width = local_component.shape
            patch_values = np.full((patch_height, patch_width, 3), 180.0, dtype=np.float32)
            return 0, 0, patch_width, patch_height, patch_values, 8.0

        with patch(
            "pdf2ppt.inpainting_engines._fit_component_background_surface",
            side_effect=fake_fit_component_background_surface,
        ):
            expanded_component, resolved_patch = inpainting_engines._resolve_component_background_patch(
                source,
                gray,
                edges,
                component,
                protected_line_mask=None,
                protected_prefill_ring_mask=None,
                kernel=kernel,
                expand_kernel=expand_kernel,
                gap_kernel=gap_kernel,
                flat_background_std_threshold=4.0,
                flat_background_edge_threshold=0.01,
                context_dilate_px=8,
                smooth_gradient_edge_threshold=0.015,
                smooth_gradient_residual_threshold=16.0,
                smooth_gradient_color_bias_max_delta=0.0,
                smooth_gradient_color_bias_residual_scale=4.0,
            )

        self.assertIsNotNone(resolved_patch)
        self.assertGreater(np.count_nonzero(expanded_component), 0)
        self.assertEqual(call_models, ["quadratic"])

    def test_resolve_component_background_patch_aligns_color_for_relaxed_quadratic(self) -> None:
        width, height = 420, 160
        source = np.full((height, width, 3), 220, dtype="uint8")
        component = np.zeros((height, width), dtype="uint8")
        component[70:118, 40:260] = 1
        gray = cv2.cvtColor(source, cv2.COLOR_BGR2GRAY)
        edges = np.zeros((height, width), dtype="uint8")
        kernel = np.ones((17, 17), dtype=np.uint8)
        expand_kernel = np.ones((5, 5), dtype=np.uint8)
        gap_kernel = np.ones((13, 13), dtype=np.uint8)

        expanded = cv2.dilate(component, expand_kernel, iterations=1)
        context_component = cv2.dilate(expanded, gap_kernel, iterations=1)
        ring_mask = cv2.dilate(context_component, kernel, iterations=1).astype(bool) & (~context_component.astype(bool))
        ring_points = np.argwhere(ring_mask)
        edge_count = max(1, int(round(float(np.count_nonzero(ring_mask)) * 0.02)))
        edges[tuple(ring_points[:edge_count].T)] = 255

        def fake_fit_component_background_surface(
            local_source: np.ndarray,
            local_component: np.ndarray,
            local_ring_mask: np.ndarray,
            *,
            context_dilate_px: int,
            model: str,
        ) -> tuple[int, int, int, int, np.ndarray, float] | None:
            patch_height, patch_width = local_component.shape
            patch_values = np.full((patch_height, patch_width, 3), 190.0, dtype=np.float32)
            return 0, 0, patch_width, patch_height, patch_values, 8.0

        with patch(
            "pdf2ppt.inpainting_engines._fit_component_background_surface",
            side_effect=fake_fit_component_background_surface,
        ):
            _, resolved_patch = inpainting_engines._resolve_component_background_patch(
                source,
                gray,
                edges,
                component,
                protected_line_mask=None,
                protected_prefill_ring_mask=None,
                kernel=kernel,
                expand_kernel=expand_kernel,
                gap_kernel=gap_kernel,
                flat_background_std_threshold=4.0,
                flat_background_edge_threshold=0.01,
                context_dilate_px=8,
                smooth_gradient_edge_threshold=0.015,
                smooth_gradient_residual_threshold=16.0,
                smooth_gradient_color_bias_max_delta=0.0,
                smooth_gradient_color_bias_residual_scale=4.0,
            )

        self.assertIsNotNone(resolved_patch)
        patch_values = resolved_patch[4]
        self.assertGreater(float(np.mean(patch_values)), 190.0)

    def test_opencv_fast_groups_nearby_small_components_before_telea(self) -> None:
        width, height = 220, 160
        base = np.full((height, width, 3), 160, dtype="uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[40:52, 50:62] = 255
        mask[40:52, 70:82] = 255

        captured_radii: list[float] = []

        def fake_inpaint(local_source: np.ndarray, local_mask: np.ndarray, radius: float, method: int) -> np.ndarray:
            captured_radii.append(radius)
            return local_source.copy()

        with (
            patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base.copy(), mask.copy())),
            patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=fake_inpaint),
        ):
            OpenCvFastInpaintingEngine(
                telea_small_component_max_span_px=24,
                telea_group_proximity_px=12,
            ).inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        self.assertEqual(len(captured_radii), 1)

    def test_opencv_fast_uses_adaptive_group_proximity_for_larger_small_components(self) -> None:
        width, height = 260, 180
        base = np.full((height, width, 3), 160, dtype="uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[30:42, 24:36] = 255
        mask[30:42, 52:64] = 255
        mask[90:122, 140:172] = 255
        mask[90:122, 188:220] = 255

        captured_radii: list[float] = []

        def fake_inpaint(local_source: np.ndarray, local_mask: np.ndarray, radius: float, method: int) -> np.ndarray:
            captured_radii.append(radius)
            return local_source.copy()

        with (
            patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base.copy(), mask.copy())),
            patch("pdf2ppt.inpainting_engines.cv2.inpaint", side_effect=fake_inpaint),
        ):
            OpenCvFastInpaintingEngine(
                telea_small_component_max_span_px=40,
                telea_group_proximity_px=20,
            ).inpaint(
                Image.fromarray(base, mode="RGB"),
                Image.fromarray(mask, mode="L"),
            )

        self.assertEqual(len(captured_radii), 3)

    def test_opencv_fast_restores_table_lines_after_masked_text_repair(self) -> None:
        width, height = 240, 180
        base_image = Image.new("RGB", (width, height), color=(248, 248, 248))
        draw = ImageDraw.Draw(base_image)
        for x in (40, 120, 200):
            draw.line((x, 20, x, 160), fill=(120, 120, 120), width=2)
        for y in (20, 70, 120, 160):
            draw.line((40, y, 200, y), fill=(120, 120, 120), width=2)
        draw.rectangle((52, 36, 108, 56), fill=(60, 60, 60))
        draw.rectangle((132, 86, 188, 106), fill=(60, 60, 60))

        mask_image = Image.new("L", (width, height), 0)
        mask_draw = ImageDraw.Draw(mask_image)
        mask_draw.rectangle((48, 30, 112, 62), fill=255)
        mask_draw.rectangle((128, 80, 192, 112), fill=255)

        base_bgr = cv2.cvtColor(np.array(base_image), cv2.COLOR_RGB2BGR)
        mask_array = np.array(mask_image, dtype=np.uint8)
        with patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base_bgr.copy(), mask_array.copy())):
            result = OpenCvFastInpaintingEngine().inpaint(base_image, mask_image)

        result_array = np.array(result, dtype=np.int16)
        base_array = np.array(base_image, dtype=np.int16)
        vertical_line_error = np.abs(result_array[30:62, 120] - base_array[30:62, 120]).mean()
        horizontal_line_error = np.abs(result_array[120, 128:192] - base_array[120, 128:192]).mean()
        self.assertLess(vertical_line_error, 15.0)
        self.assertLess(horizontal_line_error, 15.0)

    def test_opencv_fast_does_not_restore_isolated_non_grid_line(self) -> None:
        width, height = 240, 160
        base_image = Image.new("RGB", (width, height), color=(248, 248, 248))
        draw = ImageDraw.Draw(base_image)
        draw.line((30, 80, 210, 80), fill=(120, 120, 120), width=2)
        draw.rectangle((80, 62, 160, 98), fill=(60, 60, 60))

        mask_image = Image.new("L", (width, height), 0)
        mask_draw = ImageDraw.Draw(mask_image)
        mask_draw.rectangle((74, 58, 166, 102), fill=255)

        base_bgr = cv2.cvtColor(np.array(base_image), cv2.COLOR_RGB2BGR)
        mask_array = np.array(mask_image, dtype=np.uint8)
        with patch("pdf2ppt.inpainting_engines._prefill_low_texture_regions", return_value=(base_bgr.copy(), mask_array.copy())):
            result = OpenCvFastInpaintingEngine().inpaint(base_image, mask_image)

        result_array = np.array(result, dtype=np.int16)
        base_array = np.array(base_image, dtype=np.int16)
        isolated_line_error = np.abs(result_array[80, 74:166] - base_array[80, 74:166]).mean()
        self.assertGreater(isolated_line_error, 25.0)

    def test_render_overlay_background_appends_opencv_fast_telea_debug_note(self) -> None:
        image = Image.new("RGB", (60, 40), color=(35, 45, 55))
        blocks = [
            TextBlock(
                id="ocr_debug",
                source="ocr",
                bbox=(20, 12, 35, 24),
                text="demo",
                confidence=0.9,
                image_bbox=(20, 12, 35, 24),
            )
        ]
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="auto",
            inpaint_padding_px=0,
            inpaint_max_area_ratio=0.2,
        )

        class FakeOpenCvEngine(OpenCvFastInpaintingEngine):
            def inpaint(self, page_image: Image.Image, mask_image: Image.Image) -> Image.Image:
                self._last_debug_note = (
                    "Residual Telea groups: 2; group size 1-2; adaptive proximity 10-14 px; "
                    "edge density 0.010-0.080; final radius 2.10-4.30."
                )
                return page_image

        with patch(
            "pdf2ppt.inpainting_overlay.resolve_background_inpainting_engine",
            return_value=(FakeOpenCvEngine(), "opencv-fast selected"),
        ):
            result = render_overlay_background(image, blocks, fitz.Rect(0, 0, 60, 40), options=options)

        self.assertIn("edge density", result.note or "")
        self.assertIn("group size", result.note or "")
        self.assertIn("final radius", result.note or "")

    def test_opencv_fast_cleans_tight_mask_text_halo_on_smooth_background(self) -> None:
        width, height = 320, 160
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        base = np.stack(
            [
                166.0 + 0.02 * grid_x + 0.03 * grid_y,
                118.0 + 0.01 * grid_x + 0.02 * grid_y,
                182.0 + 0.015 * grid_x + 0.015 * grid_y,
            ],
            axis=2,
        ).clip(0, 255).astype("uint8")
        image = Image.fromarray(base, mode="RGB")
        draw = ImageDraw.Draw(image)
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
        text = "File-Back"
        bbox = draw.textbbox((48, 52), text, font=font)
        draw.text((48, 52), text, fill=(247, 237, 248), font=font)

        mask = Image.new("L", (width, height), 0)
        mask_draw = ImageDraw.Draw(mask)
        mask_draw.rectangle((bbox[0] + 2, bbox[1] + 2, bbox[2] - 2, bbox[3] - 2), fill=255)

        result = OpenCvFastInpaintingEngine().inpaint(image, mask)

        x0 = max(0, bbox[0] - 3)
        y0 = max(0, bbox[1] - 3)
        x1 = min(width, bbox[2] + 3)
        y1 = min(height, bbox[3] + 3)
        result_array = np.array(result, dtype=np.int16)
        base_array = base.astype(np.int16)
        mae = np.abs(result_array[y0:y1, x0:x1] - base_array[y0:y1, x0:x1]).mean()
        self.assertLess(mae, 1.0)

    def test_opencv_fast_restores_curved_gradient_background(self) -> None:
        width, height = 260, 160
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        channel_r = 90.0 + 0.55 * grid_x + 0.18 * grid_y + 0.0022 * ((grid_x - 140.0) ** 2)
        channel_r += 0.0008 * ((grid_y - 80.0) ** 2)
        channel_g = 70.0 + 0.32 * grid_x + 0.24 * grid_y + 0.0015 * ((grid_x - 120.0) ** 2)
        channel_b = 120.0 + 0.18 * grid_x + 0.16 * grid_y + 0.0018 * ((grid_x - 150.0) ** 2)
        channel_b += 0.0009 * grid_x * grid_y / 10.0
        base = np.stack(
            [channel_r, channel_g, channel_b],
            axis=2,
        ).clip(0, 255).astype("uint8")
        mask = np.zeros((height, width), dtype="uint8")
        mask[48:118, 70:210] = 255

        result = OpenCvFastInpaintingEngine().inpaint(
            Image.fromarray(base, mode="RGB"),
            Image.fromarray(mask, mode="L"),
        )

        result_array = np.array(result, dtype=np.int16)
        base_array = base.astype(np.int16)
        mae = np.abs(result_array - base_array)[mask > 0].mean()
        self.assertLess(mae, 1.5)

    def test_targeted_file_back_lab_correction_improves_local_ring_alignment(self) -> None:
        width, height = 260, 160
        grid_y, grid_x = np.indices((height, width), dtype=np.float32)
        base = np.stack(
            [
                195.0 + 0.14 * grid_x + 0.08 * grid_y,
                170.0 + 0.10 * grid_x + 0.07 * grid_y,
                205.0 + 0.12 * grid_x + 0.05 * grid_y,
            ],
            axis=2,
        ).clip(0, 255).astype("uint8")
        page_image = Image.fromarray(base, mode="RGB")
        repaired = base.astype(np.float32)
        repaired_region = repaired[52:108, 78:198]
        repaired_region_mean = repaired_region.mean(axis=(0, 1), keepdims=True)
        repaired[52:108, 78:198] = repaired_region_mean + (repaired_region - repaired_region_mean) * 0.45
        repaired[52:108, 78:198, 0] -= 20.0
        repaired[52:108, 78:198, 1] -= 14.0
        repaired[52:108, 78:198, 2] -= 18.0
        repaired = np.clip(repaired, 0, 255).astype("uint8")
        repaired_image = Image.fromarray(repaired, mode="RGB")
        block = TextBlock(
            id="ocr_file_back",
            source="ocr",
            bbox=(78, 52, 198, 108),
            text="歸檔(File-Back)",
            confidence=0.95,
            image_bbox=(78, 52, 198, 108),
        )
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=0,
        )

        corrected_image, debug_images, note = _apply_targeted_file_back_color_correction(
            page_image,
            repaired_image,
            [block],
            fitz.Rect(0, 0, width, height),
            options=options,
        )

        repaired_arr = np.array(repaired_image, dtype=np.uint8)
        corrected_arr = np.array(corrected_image, dtype=np.uint8)
        component = np.zeros((height, width), dtype=np.uint8)
        component[52:108, 78:198] = 1
        inner = cv2.dilate(component, np.ones((13, 13), np.uint8), iterations=1).astype(bool)
        outer = cv2.dilate(inner.astype(np.uint8), np.ones((21, 21), np.uint8), iterations=1).astype(bool)
        ring = outer & (~inner)
        region = component.astype(bool)
        ring_mean = base[ring].astype(np.float32).mean(axis=0)
        repaired_gap = np.abs(repaired_arr[region].astype(np.float32).mean(axis=0) - ring_mean).mean()
        corrected_gap = np.abs(corrected_arr[region].astype(np.float32).mean(axis=0) - ring_mean).mean()
        repaired_lab = cv2.cvtColor(repaired_arr, cv2.COLOR_RGB2LAB).astype(np.float32)
        corrected_lab = cv2.cvtColor(corrected_arr, cv2.COLOR_RGB2LAB).astype(np.float32)
        repaired_span = np.percentile(repaired_lab[:, :, 0][region], 90) - np.percentile(
            repaired_lab[:, :, 0][region], 10
        )
        corrected_span = np.percentile(corrected_lab[:, :, 0][region], 90) - np.percentile(
            corrected_lab[:, :, 0][region], 10
        )

        self.assertLess(corrected_gap, repaired_gap)
        self.assertGreater(corrected_span, repaired_span)
        self.assertIn("File-Back", note or "")
        self.assertIn("file_back_corrected", debug_images)
        self.assertIn("file_back_mask", debug_images)

    def test_apply_targeted_footer_label_color_correction_reduces_mask_ring_gap(self) -> None:
        width, height = 320, 180
        base = np.full((height, width, 3), 238, dtype=np.uint8)
        mask = np.zeros((height, width), dtype=np.uint8)
        cv2.fillConvexPoly(mask, np.array([[68, 138], [248, 136], [246, 176], [70, 174]], dtype=np.int32), 1)
        inner = cv2.dilate(mask, np.ones((5, 5), np.uint8), iterations=1).astype(bool)
        near_outer = cv2.dilate(inner.astype(np.uint8), np.ones((17, 17), np.uint8), iterations=1).astype(bool)
        near_ring = near_outer & (~inner)
        base_with_halo = base.copy().astype(np.float32)
        base_with_halo[near_ring] -= 14.0
        base_with_halo = np.clip(base_with_halo, 0, 255).astype(np.uint8)
        repaired = base_with_halo.copy().astype(np.float32)
        repaired[mask.astype(bool)] -= 24.0
        repaired = np.clip(repaired, 0, 255).astype(np.uint8)
        page_image = Image.fromarray(base_with_halo, mode="RGB")
        repaired_image = Image.fromarray(repaired, mode="RGB")
        block = TextBlock(
            id="ocr_footer_label",
            source="ocr",
            bbox=(68.0, 136.0, 248.0, 176.0),
            text="同義改寫 (100%)",
            confidence=0.95,
            block_role="body",
            image_bbox=(68.0, 136.0, 248.0, 176.0),
            image_polygon=((68.0, 138.0), (248.0, 136.0), (246.0, 176.0), (70.0, 174.0)),
        )
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=0,
        )

        corrected_image, debug_images, note = _apply_targeted_footer_label_color_correction(
            page_image,
            repaired_image,
            [block],
            fitz.Rect(0, 0, width, height),
            options=options,
        )

        outer = cv2.dilate(inner.astype(np.uint8), np.ones((17, 17), np.uint8), iterations=1).astype(bool)
        ring = outer & (~inner)
        broad_outer = cv2.dilate(inner.astype(np.uint8), np.ones((49, 49), np.uint8), iterations=1).astype(bool)
        broad_ring = broad_outer & (~outer)
        repaired_arr = np.array(repaired_image, dtype=np.uint8)
        corrected_arr = np.array(corrected_image, dtype=np.uint8)
        broad_mean = base_with_halo[broad_ring].astype(np.float32).mean(axis=0)
        repaired_broad_gap = np.abs(repaired_arr[mask.astype(bool)].astype(np.float32).mean(axis=0) - broad_mean).mean()
        corrected_broad_gap = np.abs(corrected_arr[mask.astype(bool)].astype(np.float32).mean(axis=0) - broad_mean).mean()

        self.assertLess(corrected_broad_gap, repaired_broad_gap)
        self.assertIn("broad lift", (note or "").lower())
        self.assertIn("footer label flat tone correction", (note or "").lower())
        self.assertIn("footer_label_01_corrected", debug_images)
        self.assertIn("footer_label_01_mask", debug_images)

    def test_apply_targeted_footer_label_color_correction_uses_shared_footer_target(self) -> None:
        width, height = 520, 180
        base = np.full((height, width, 3), 238, dtype=np.uint8)
        left_mask = np.zeros((height, width), dtype=np.uint8)
        right_mask = np.zeros((height, width), dtype=np.uint8)
        cv2.fillConvexPoly(left_mask, np.array([[28, 138], [208, 136], [206, 176], [30, 174]], dtype=np.int32), 1)
        cv2.fillConvexPoly(right_mask, np.array([[288, 138], [468, 136], [466, 176], [290, 174]], dtype=np.int32), 1)
        union_mask = (left_mask | right_mask).astype(np.uint8)
        inner = cv2.dilate(union_mask, np.ones((5, 5), np.uint8), iterations=1).astype(bool)
        near_outer = cv2.dilate(inner.astype(np.uint8), np.ones((17, 17), np.uint8), iterations=1).astype(bool)
        base_with_halo = base.copy().astype(np.float32)
        base_with_halo[:, :260] -= np.array([18.0, 18.0, 18.0], dtype=np.float32)
        base_with_halo = np.clip(base_with_halo, 0, 255).astype(np.uint8)
        repaired = base_with_halo.copy().astype(np.float32)
        repaired[left_mask.astype(bool)] -= 18.0
        repaired[right_mask.astype(bool)] -= 18.0
        repaired = np.clip(repaired, 0, 255).astype(np.uint8)
        page_image = Image.fromarray(base_with_halo, mode="RGB")
        repaired_image = Image.fromarray(repaired, mode="RGB")
        left_block = TextBlock(
            id="ocr_footer_label_left",
            source="ocr",
            bbox=(28.0, 136.0, 208.0, 176.0),
            text="高頻FAQ(98.8%)",
            confidence=0.95,
            block_role="body",
            image_bbox=(28.0, 136.0, 208.0, 176.0),
            image_polygon=((28.0, 138.0), (208.0, 136.0), (206.0, 176.0), (30.0, 174.0)),
        )
        right_block = TextBlock(
            id="ocr_footer_label_right",
            source="ocr",
            bbox=(288.0, 136.0, 468.0, 176.0),
            text="邊界題 (86.7%)",
            confidence=0.95,
            block_role="body",
            image_bbox=(288.0, 136.0, 468.0, 176.0),
            image_polygon=((288.0, 138.0), (468.0, 136.0), (466.0, 176.0), (290.0, 174.0)),
        )
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            inpaint_engine="opencv-fast",
            inpaint_padding_px=0,
        )

        corrected_image, _debug_images, note = _apply_targeted_footer_label_color_correction(
            page_image,
            repaired_image,
            [left_block, right_block],
            fitz.Rect(0, 0, width, height),
            options=options,
        )

        corrected_arr = np.array(corrected_image, dtype=np.uint8)
        left_mean = corrected_arr[left_mask.astype(bool)].astype(np.float32).mean(axis=0)
        right_mean = corrected_arr[right_mask.astype(bool)].astype(np.float32).mean(axis=0)

        self.assertLess(np.abs(left_mean - right_mean).mean(), 1.0)
        self.assertIn("footer label flat tone correction", (note or "").lower())

    def test_apply_targeted_footer_label_color_correction_prefers_bright_background_percentile(self) -> None:
        width, height = 320, 180
        base = np.full((height, width, 3), 238, dtype=np.uint8)
        mask = np.zeros((height, width), dtype=np.uint8)
        cv2.fillConvexPoly(mask, np.array([[68, 138], [248, 136], [246, 176], [70, 174]], dtype=np.int32), 1)
        inner = cv2.dilate(mask, np.ones((5, 5), np.uint8), iterations=1).astype(bool)
        near_outer = cv2.dilate(inner.astype(np.uint8), np.ones((17, 17), np.uint8), iterations=1).astype(bool)
        broad_outer = cv2.dilate(inner.astype(np.uint8), np.ones((49, 49), np.uint8), iterations=1).astype(bool)
        broad_ring = broad_outer & (~near_outer)
        page = base.copy().astype(np.float32)
        page[near_outer & (~inner)] -= 16.0
        bright_ring = broad_ring.copy()
        bright_ring[:, :80] = False
        page[broad_ring] = 228.0
        page[bright_ring] = 242.0
        page = np.clip(page, 0, 255).astype(np.uint8)
        repaired = page.copy().astype(np.float32)
        repaired[mask.astype(bool)] -= 20.0
        repaired = np.clip(repaired, 0, 255).astype(np.uint8)

        corrected_image, _debug_images, _note = _apply_targeted_footer_label_color_correction(
            Image.fromarray(page, mode="RGB"),
            Image.fromarray(repaired, mode="RGB"),
            [
                TextBlock(
                    id="ocr_footer_label",
                    source="ocr",
                    bbox=(68.0, 136.0, 248.0, 176.0),
                    text="同義改寫 (100%)",
                    confidence=0.95,
                    block_role="body",
                    image_bbox=(68.0, 136.0, 248.0, 176.0),
                    image_polygon=((68.0, 138.0), (248.0, 136.0), (246.0, 176.0), (70.0, 174.0)),
                )
            ],
            fitz.Rect(0, 0, width, height),
            options=ConversionOptions(
                input_path=Path("input.pdf"),
                output_path=Path("output.pptx"),
                report_path=Path("output.report.json"),
                inpaint_engine="opencv-fast",
                inpaint_padding_px=0,
            ),
        )

        corrected_arr = np.array(corrected_image, dtype=np.uint8)
        corrected_mean = corrected_arr[mask.astype(bool)].astype(np.float32).mean(axis=0)
        broad_mean = page[broad_ring].astype(np.float32).mean(axis=0)
        broad_p85 = np.percentile(page[broad_ring].astype(np.float32), 85, axis=0)

        self.assertLess(
            np.abs(corrected_mean - broad_p85).mean(),
            np.abs(corrected_mean - broad_mean).mean(),
        )

class AnalyzePagePerformanceTests(unittest.TestCase):
    @patch("pdf2ppt.pipeline.extract_image_elements", return_value=[])
    @patch("pdf2ppt.pipeline.render_page_image")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("elements", None))
    @patch(
        "pdf2ppt.pipeline.score_page",
        return_value=QualityScore(
            text_confidence=0.98,
            layout_overlap_score=0.95,
            style_recovery_score=1.0,
            editable_ratio=1.0,
        ),
    )
    @patch("pdf2ppt.pipeline.select_text_blocks")
    @patch("pdf2ppt.pipeline.classify_page", return_value="digital")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=100,
            native_text_area_ratio=0.2,
            image_area_ratio=0.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks")
    def test_analyze_page_skips_render_for_elements_mode(
        self,
        extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        select_text_blocks_mock: unittest.mock.Mock,
        _score_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        render_page_image_mock: unittest.mock.Mock,
        extract_image_elements_mock: unittest.mock.Mock,
    ) -> None:
        native_blocks = [
            TextBlock(
                id="n1",
                source="native",
                bbox=(10, 10, 100, 30),
                text="Hello",
                confidence=0.99,
            )
        ]
        extract_native_text_blocks_mock.return_value = (native_blocks, [])
        select_text_blocks_mock.return_value = native_blocks
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
        )

        result = analyze_page(page, options, ocr_engine=unittest.mock.Mock())

        self.assertEqual(result.background_mode, "elements")
        render_page_image_mock.assert_not_called()
        extract_image_elements_mock.assert_called_once()

    @patch("pdf2ppt.pipeline.render_overlay_background")
    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"jpeg-bytes")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("overlay", None))
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_prefers_approved_boxes_over_raw_ocr(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        _encode_background_mock: unittest.mock.Mock,
        render_overlay_background_mock: unittest.mock.Mock,
    ) -> None:
        preview_image = Image.new("RGB", (200, 100), (10, 10, 10))
        render_overlay_background_mock.return_value = SimpleNamespace(
            image=preview_image,
            engine_name="opencv-fast",
            note="ok",
            mask_image=None,
            debug_images={},
        )
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            approved_ocr_blocks_by_page={
                1: [
                    TextBlock(
                        id="approved_1",
                        source="ocr",
                        bbox=(10, 20, 50, 60),
                        text="Approved",
                        confidence=1.0,
                        image_bbox=(10, 20, 50, 60),
                        image_polygon=((10, 20), (50, 20), (50, 60), (10, 60)),
                    )
                ]
            },
            approved_ocr_image_size_by_page={1: (200, 100)},
        )
        ocr_engine = unittest.mock.Mock()

        with patch("pdf2ppt.pipeline.render_page_image", return_value=preview_image):
            result = analyze_page(page, options, ocr_engine=ocr_engine)

        ocr_engine.extract_text_blocks.assert_not_called()
        self.assertEqual([block.text for block in result.text_blocks], ["Approved"])
        self.assertEqual(result.background_inpaint_engine, "opencv-fast")

    @patch("pdf2ppt.pipeline.render_overlay_background")
    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"jpeg-bytes")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("overlay", None))
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_recognizes_text_for_blank_approved_box(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        _encode_background_mock: unittest.mock.Mock,
        render_overlay_background_mock: unittest.mock.Mock,
    ) -> None:
        preview_image = Image.new("RGB", (200, 100), (10, 10, 10))
        render_overlay_background_mock.return_value = SimpleNamespace(
            image=preview_image,
            engine_name="opencv-fast",
            note="ok",
            mask_image=None,
            debug_images={},
        )
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            approved_ocr_blocks_by_page={
                1: [
                    TextBlock(
                        id="approved_blank_1",
                        source="ocr",
                        bbox=(10, 20, 50, 60),
                        text="",
                        confidence=1.0,
                        image_bbox=(10, 20, 50, 60),
                    )
                ]
            },
            approved_ocr_image_size_by_page={1: (200, 100)},
        )
        ocr_engine = unittest.mock.Mock()
        ocr_engine.recognize_text_in_box.return_value = ("Manual OCR", 0.88)

        with patch("pdf2ppt.pipeline.render_page_image", return_value=preview_image):
            result = analyze_page(page, options, ocr_engine=ocr_engine)

        ocr_engine.extract_text_blocks.assert_not_called()
        ocr_engine.recognize_text_in_box.assert_called_once()
        self.assertEqual([block.text for block in result.text_blocks], ["Manual OCR"])

    @patch("pdf2ppt.pipeline.render_overlay_background")
    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"png-bytes")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("overlay", None))
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_encodes_overlay_background_as_png(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        encode_background_mock: unittest.mock.Mock,
        render_overlay_background_mock: unittest.mock.Mock,
    ) -> None:
        preview_image = Image.new("RGB", (200, 100), (10, 10, 10))
        render_overlay_background_mock.return_value = SimpleNamespace(
            image=preview_image,
            engine_name="opencv-fast",
            note="ok",
            mask_image=None,
            debug_images={},
        )
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            background_image_format="jpeg",
            background_jpeg_quality=70,
        )
        ocr_engine = unittest.mock.Mock()
        ocr_engine.extract_text_blocks.return_value = SimpleNamespace(blocks=[], image=preview_image)

        with patch("pdf2ppt.pipeline.render_page_image", return_value=preview_image):
            result = analyze_page(page, options, ocr_engine=ocr_engine)

        self.assertEqual(result.background_image_bytes, b"png-bytes")
        encode_background_mock.assert_called_once_with(
            unittest.mock.ANY,
            image_format="png",
            jpeg_quality=70,
        )
        encoded_background = encode_background_mock.call_args.args[0]
        self.assertEqual(encoded_background.size, (200, 100))

    @patch("pdf2ppt.pipeline.render_overlay_background")
    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"png-bytes")
    @patch(
        "pdf2ppt.pipeline.select_text_blocks",
        return_value=[
            TextBlock(
                id="ocr_overlay_1",
                source="ocr",
                bbox=(10, 10, 80, 30),
                text="overlay",
                confidence=0.9,
            )
        ],
    )
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("overlay", None))
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_uses_render_dpi_background_for_overlay(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _select_text_blocks_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        encode_background_mock: unittest.mock.Mock,
        render_overlay_background_mock: unittest.mock.Mock,
    ) -> None:
        ocr_image = Image.new("RGB", (200, 100), (10, 10, 10))
        render_overlay_background_mock.return_value = SimpleNamespace(
            image=ocr_image,
            engine_name="opencv-fast",
            note="ok",
            mask_image=None,
            debug_images={},
        )
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            render_dpi=144,
            background_dpi=96,
            background_image_format="jpeg",
        )
        ocr_engine = unittest.mock.Mock()
        ocr_engine.extract_text_blocks.return_value = SimpleNamespace(blocks=[], image=ocr_image)

        with patch("pdf2ppt.pipeline.render_page_image", return_value=ocr_image) as render_page_image_mock:
            analyze_page(page, options, ocr_engine=ocr_engine)

        render_page_image_mock.assert_called_once_with(page, dpi=144)
        background_image = render_overlay_background_mock.call_args.args[0]
        self.assertEqual(background_image.size, (200, 100))
        encode_background_mock.assert_called_once_with(
            ocr_image,
            image_format="png",
            jpeg_quality=82,
        )

    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"jpeg-bytes")
    @patch("pdf2ppt.pipeline.render_page_image")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("full-page", "fallback"))
    @patch(
        "pdf2ppt.pipeline.score_page",
        return_value=QualityScore(
            text_confidence=0.8,
            layout_overlap_score=0.7,
            style_recovery_score=0.45,
            editable_ratio=0.2,
        ),
    )
    @patch("pdf2ppt.pipeline.select_text_blocks", return_value=[])
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_renders_background_with_background_dpi(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _select_text_blocks_mock: unittest.mock.Mock,
        _score_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        render_page_image_mock: unittest.mock.Mock,
        encode_background_mock: unittest.mock.Mock,
    ) -> None:
        ocr_image = Image.new("RGB", (200, 100), (10, 10, 10))
        render_page_image_mock.return_value = ocr_image
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            render_dpi=144,
            background_dpi=96,
            background_image_format="jpeg",
            background_jpeg_quality=70,
        )
        ocr_engine = unittest.mock.Mock()
        ocr_engine.extract_text_blocks.return_value = SimpleNamespace(blocks=[], image=ocr_image)

        result = analyze_page(page, options, ocr_engine=ocr_engine)

        self.assertEqual(result.background_image_bytes, b"jpeg-bytes")
        self.assertEqual(render_page_image_mock.call_args_list[0].kwargs["dpi"], 144)
        self.assertEqual(render_page_image_mock.call_count, 1)
        encode_background_mock.assert_called_once_with(
            unittest.mock.ANY,
            image_format="jpeg",
            jpeg_quality=70,
        )
        resized_background = encode_background_mock.call_args.args[0]
        self.assertEqual(resized_background.size, (133, 67))

    @patch("pdf2ppt.pipeline.pil_to_image_bytes", return_value=b"jpeg-bytes")
    @patch("pdf2ppt.pipeline.render_page_image")
    @patch("pdf2ppt.pipeline.choose_background_mode", return_value=("full-page", "fallback"))
    @patch(
        "pdf2ppt.pipeline.score_page",
        return_value=QualityScore(
            text_confidence=0.8,
            layout_overlap_score=0.7,
            style_recovery_score=0.45,
            editable_ratio=0.2,
        ),
    )
    @patch("pdf2ppt.pipeline.select_text_blocks", return_value=[])
    @patch("pdf2ppt.pipeline.classify_page", return_value="scanned")
    @patch(
        "pdf2ppt.pipeline.compute_page_signals",
        return_value=PageSignals(
            native_char_count=0,
            native_text_area_ratio=0.0,
            image_area_ratio=1.0,
            drawing_count=0,
        ),
    )
    @patch("pdf2ppt.pipeline.extract_native_text_blocks", return_value=([], []))
    def test_analyze_page_rerenders_background_when_doc_orientation_enabled(
        self,
        _extract_native_text_blocks_mock: unittest.mock.Mock,
        _compute_page_signals_mock: unittest.mock.Mock,
        _classify_page_mock: unittest.mock.Mock,
        _select_text_blocks_mock: unittest.mock.Mock,
        _score_page_mock: unittest.mock.Mock,
        _choose_background_mode_mock: unittest.mock.Mock,
        render_page_image_mock: unittest.mock.Mock,
        encode_background_mock: unittest.mock.Mock,
    ) -> None:
        ocr_image = Image.new("RGB", (200, 100), (10, 10, 10))
        background_image = Image.new("RGB", (120, 60), (20, 20, 20))
        render_page_image_mock.side_effect = [ocr_image, background_image]
        page = SimpleNamespace(number=0, rect=fitz.Rect(0, 0, 320, 240))
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            render_dpi=144,
            background_dpi=96,
            background_image_format="jpeg",
            background_jpeg_quality=70,
            ocr_use_doc_orientation=True,
        )
        ocr_engine = unittest.mock.Mock()
        ocr_engine.extract_text_blocks.return_value = SimpleNamespace(blocks=[], image=ocr_image)

        result = analyze_page(page, options, ocr_engine=ocr_engine)

        self.assertEqual(result.background_image_bytes, b"jpeg-bytes")
        self.assertEqual(render_page_image_mock.call_args_list[0].kwargs["dpi"], 144)
        self.assertEqual(render_page_image_mock.call_args_list[1].kwargs["dpi"], 96)
        encode_background_mock.assert_called_once_with(
            background_image,
            image_format="jpeg",
            jpeg_quality=70,
        )


class OcrModelConfigTests(unittest.TestCase):
    def test_suppress_known_paddle_runtime_warnings_filters_only_ccache_warning(self) -> None:
        with warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always")
            with suppress_known_paddle_runtime_warnings():
                warnings.warn(
                    "No ccache found. Please be aware that recompiling all source files may be required.",
                    UserWarning,
                )
                warnings.warn("keep this warning", UserWarning)

        self.assertEqual(len(caught), 1)
        self.assertEqual(str(caught[0].message), "keep this warning")

    def test_build_local_ocr_model_kwargs_uses_repo_local_model_root(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            model_root = Path(tmp_dir) / "model"
            kwargs = build_local_ocr_model_kwargs(
                model_root=model_root,
                lang="ch",
                use_doc_orientation=False,
                use_doc_unwarping=False,
                use_textline_orientation=False,
            )

            self.assertEqual(kwargs["text_detection_model_name"], "PP-OCRv5_server_det")
            self.assertEqual(kwargs["text_recognition_model_name"], "PP-OCRv5_server_rec")
            self.assertEqual(Path(kwargs["text_detection_model_dir"]).parent, model_root.resolve())
            self.assertNotIn("doc_orientation_classify_model_dir", kwargs)
            self.assertNotIn("textline_orientation_model_dir", kwargs)

    def test_build_local_ocr_model_kwargs_can_enable_doc_orientation(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            model_root = Path(tmp_dir) / "model"
            kwargs = build_local_ocr_model_kwargs(
                model_root=model_root,
                lang="ch",
                use_doc_orientation=True,
                use_doc_unwarping=False,
                use_textline_orientation=False,
            )

            self.assertEqual(kwargs["doc_orientation_classify_model_name"], "PP-LCNet_x1_0_doc_ori")
            self.assertEqual(Path(kwargs["doc_orientation_classify_model_dir"]).parent, model_root.resolve())

    def test_ensure_local_model_dir_reuses_existing_paddlex_cache(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            temp_root = Path(tmp_dir)
            cache_root = temp_root / "cache"
            cached_model_dir = cache_root / "PP-OCRv5_server_det"
            cached_model_dir.mkdir(parents=True)
            (cached_model_dir / "inference.json").write_text("{}", encoding="utf-8")

            with patch("pdf2ppt.ocr.PADDLEX_OFFICIAL_MODEL_CACHE_DIR", cache_root):
                local_model_dir = ensure_local_model_dir(temp_root / "model", "PP-OCRv5_server_det")

            self.assertTrue(local_model_dir.exists())
            self.assertTrue((local_model_dir / "inference.json").exists())


class CliTests(unittest.TestCase):
    def test_doc_unwarping_disabled_by_default(self) -> None:
        parser = build_parser()
        args = parser.parse_args(["input.pdf", "output.pptx"])
        self.assertFalse(args.enable_doc_unwarping)
        self.assertFalse(args.enable_doc_orientation)
        self.assertEqual(args.ocr_model_root, Path("model"))
        self.assertFalse(args.enable_textline_orientation)
        self.assertEqual(args.inpaint_engine, "opencv-fast")
        self.assertIsNone(args.ocr_det_thresh)
        self.assertIsNone(args.ocr_det_box_thresh)
        self.assertIsNone(args.ocr_drop_score)
        self.assertEqual(args.ocr_batch_size, 3)
        self.assertEqual(args.inpaint_padding_px, 6)
        self.assertAlmostEqual(args.inpaint_max_area_ratio, 0.12)
        self.assertEqual(args.background_dpi, 110)
        self.assertEqual(args.background_format, "jpeg")
        self.assertEqual(args.background_jpeg_quality, 82)
        self.assertEqual(args.log_level, "INFO")

    def test_doc_unwarping_can_be_enabled(self) -> None:
        parser = build_parser()
        args = parser.parse_args(["input.pdf", "output.pptx", "--enable-doc-unwarping"])
        self.assertTrue(args.enable_doc_unwarping)

    def test_orientation_flags_and_model_root_can_be_enabled(self) -> None:
        parser = build_parser()
        args = parser.parse_args(
            [
                "input.pdf",
                "output.pptx",
                "--ocr-model-root",
                "local-models",
                "--enable-doc-orientation",
                "--enable-textline-orientation",
            ]
        )
        self.assertEqual(args.ocr_model_root, Path("local-models"))
        self.assertTrue(args.enable_doc_orientation)
        self.assertTrue(args.enable_textline_orientation)

    @patch("pdf2ppt.cli.convert_pdf")
    def test_main_does_not_enable_debug_dir_by_default(self, convert_pdf_mock: unittest.mock.Mock) -> None:
        convert_pdf_mock.return_value = SimpleNamespace(
            input_path="input.pdf",
            output_path="output.pptx",
            pages=[object()],
        )

        exit_code = main(["input.pdf", "output.pptx"])

        self.assertEqual(exit_code, 0)
        options = convert_pdf_mock.call_args.args[0]
        self.assertIsNone(options.debug_dir)

    def test_inpaint_flags_can_be_configured(self) -> None:
        parser = build_parser()
        args = parser.parse_args(
            [
                "input.pdf",
                "output.pptx",
                "--ocr-det-thresh",
                "0.55",
                "--ocr-det-box-thresh",
                "0.6",
                "--ocr-drop-score",
                "0.65",
                "--ocr-batch-size",
                "6",
                "--inpaint-engine",
                "opencv-fast",
                "--inpaint-padding-px",
                "10",
                "--inpaint-max-area-ratio",
                "0.2",
                "--inpaint-model-root",
                "custom-lama",
                "--inpaint-onnx-cuda-provider",
                "CUDAExecutionProvider",
                "--inpaint-onnx-execution-mode",
                "parallel",
                "--inpaint-max-side-px",
                "1024",
                "--background-dpi",
                "96",
                "--background-format",
                "png",
                "--background-jpeg-quality",
                "70",
                "--log-level",
                "DEBUG",
            ]
        )
        self.assertAlmostEqual(args.ocr_det_thresh, 0.55)
        self.assertAlmostEqual(args.ocr_det_box_thresh, 0.6)
        self.assertAlmostEqual(args.ocr_drop_score, 0.65)
        self.assertEqual(args.ocr_batch_size, 6)
        self.assertEqual(args.inpaint_engine, "opencv-fast")
        self.assertEqual(args.inpaint_padding_px, 10)
        self.assertAlmostEqual(args.inpaint_max_area_ratio, 0.2)
        self.assertEqual(args.inpaint_model_root, Path("custom-lama"))
        self.assertEqual(args.inpaint_onnx_cuda_provider, "CUDAExecutionProvider")
        self.assertEqual(args.inpaint_onnx_execution_mode, "parallel")
        self.assertEqual(args.inpaint_max_side_px, 1024)
        self.assertEqual(args.background_dpi, 96)
        self.assertEqual(args.background_format, "png")
        self.assertEqual(args.background_jpeg_quality, 70)
        self.assertEqual(args.log_level, "DEBUG")

    def test_inpaint_engine_defaults_to_opencv_fast(self) -> None:
        parser = build_parser()
        args = parser.parse_args(["input.pdf", "output.pptx"])
        self.assertEqual(args.inpaint_engine, "opencv-fast")

    @patch("pdf2ppt.cli.convert_pdf")
    def test_main_passes_custom_ocr_batch_size(self, convert_pdf_mock: unittest.mock.Mock) -> None:
        convert_pdf_mock.return_value = SimpleNamespace(
            input_path="input.pdf",
            output_path="output.pptx",
            pages=[object()],
        )

        exit_code = main(["input.pdf", "output.pptx", "--ocr-batch-size", "4"])

        self.assertEqual(exit_code, 0)
        options = convert_pdf_mock.call_args.args[0]
        self.assertEqual(options.ocr_batch_size, 4)

    def test_format_progress_line_reports_completion(self) -> None:
        line = format_progress_line(2, 4, width=8)
        self.assertEqual(line, "Converting pages [####----] 2/4 ( 50%)")

    def test_build_progress_callback_writes_tty_progress_bar(self) -> None:
        class TtyBuffer(StringIO):
            def isatty(self) -> bool:
                return True

        buffer = TtyBuffer()
        callback = build_progress_callback(buffer, width=4)
        callback(0, 2)
        callback(1, 2)
        callback(2, 2)
        self.assertIn("\rConverting pages [----] 0/2 (  0%)", buffer.getvalue())
        self.assertIn("\rConverting pages [##--] 1/2 ( 50%)", buffer.getvalue())
        self.assertTrue(buffer.getvalue().endswith("\rConverting pages [####] 2/2 (100%)\n"))


class FontSizingTests(unittest.TestCase):
    def test_pil_to_image_bytes_supports_jpeg(self) -> None:
        image = Image.new("RGB", (24, 24), (120, 140, 160))
        encoded = pil_to_image_bytes(image, image_format="jpeg", jpeg_quality=70)
        self.assertTrue(encoded.startswith(b"\xff\xd8\xff"))

    def test_resolve_background_render_dpi_uses_option(self) -> None:
        options = ConversionOptions(
            input_path=Path("input.pdf"),
            output_path=Path("output.pptx"),
            report_path=Path("output.report.json"),
            background_dpi=96,
        )
        self.assertEqual(resolve_background_render_dpi(options), 96)

    def test_default_font_family_by_script(self) -> None:
        self.assertEqual(default_font_family("latin"), "DejaVu Sans")
        self.assertEqual(default_font_family("cjk"), "Noto Sans CJK TC")

    def test_classify_text_script(self) -> None:
        self.assertEqual(classify_text_script("向量搜尋"), "cjk")
        self.assertEqual(classify_text_script("VECTOR"), "latin")
        self.assertEqual(classify_text_script("2026"), "numeric")
        self.assertEqual(classify_text_script("RAG 2026 向量"), "mixed")

    def test_estimate_font_size_grows_with_box_height(self) -> None:
        small = estimate_font_size("VECTOR", (0, 0, 120, 18))
        large = estimate_font_size("VECTOR", (0, 0, 120, 36))
        self.assertGreater(large, small)

    def test_estimate_font_size_stays_reasonable(self) -> None:
        estimate = estimate_font_size("RAG 2026：當向量搜尋遇上代理推理", (0, 0, 1200, 120))
        self.assertGreaterEqual(estimate, 6.0)
        self.assertLessEqual(estimate, 96.0)

    def test_build_text_fit_debug_entry_contains_error_metrics(self) -> None:
        block = TextBlock(
            id="ocr_debug_1",
            source="ocr",
            bbox=(0, 0, 200, 24),
            text="VECTOR EMBEDDINGS",
            confidence=0.9,
            font_size=18,
        )
        entry = build_text_fit_debug_entry(block)
        self.assertEqual(entry["id"], "ocr_debug_1")
        self.assertIn("width_error_ratio", entry)
        self.assertIn("height_error_ratio", entry)
        self.assertEqual(entry["target_bbox_pt"]["width"], 200.0)

    def test_fit_text_frame_uses_fit_text_for_ocr_blocks(self) -> None:
        calls: list[dict[str, object]] = []

        class DummyTextFrame:
            def fit_text(self, **kwargs: object) -> None:
                calls.append(kwargs)

        block = TextBlock(
            id="ocr_fit_1",
            source="ocr",
            bbox=(0, 0, 140, 24),
            text="VECTOR EMBEDDINGS",
            confidence=0.9,
            font_size=12,
        )
        self.assertTrue(fit_text_frame(DummyTextFrame(), block, scale_x=1.0, scale_y=1.0))
        self.assertEqual(len(calls), 1)
        self.assertEqual(calls[0]["font_family"], "DejaVu Sans")
        self.assertGreaterEqual(calls[0]["max_size"], 6)

    @patch("pdf2ppt.ppt_render.resolve_ocr_fit_max_size", return_value=18)
    def test_fit_text_frame_uses_resolved_cap_for_ocr_blocks(self, _resolve_max_size: unittest.mock.Mock) -> None:
        calls: list[dict[str, object]] = []

        class DummyTextFrame:
            def fit_text(self, **kwargs: object) -> None:
                calls.append(kwargs)

        block = TextBlock(
            id="ocr_fit_2",
            source="ocr",
            bbox=(0, 0, 200, 36),
            text="NotebookLM",
            confidence=0.9,
            font_size=12,
        )
        self.assertTrue(fit_text_frame(DummyTextFrame(), block, scale_x=1.0, scale_y=1.0))
        self.assertEqual(calls[0]["max_size"], 18)

    def test_should_wrap_text_block_disables_wrap_for_single_line_ocr(self) -> None:
        self.assertFalse(
            should_wrap_text_block(
                TextBlock(
                    id="ocr_wrap_1",
                    source="ocr",
                    bbox=(0, 0, 100, 20),
                    text="VECTOR EMBEDDINGS",
                    confidence=0.9,
                )
            )
        )
        self.assertTrue(
            should_wrap_text_block(
                TextBlock(
                    id="ocr_wrap_2",
                    source="ocr",
                    bbox=(0, 0, 100, 20),
                    text="LINE 1\nLINE 2",
                    confidence=0.9,
                )
            )
        )

    def test_resolve_fallback_font_size_pt_shrinks_single_line_text_to_fit_bbox(self) -> None:
        block = TextBlock(
            id="ocr_fallback_font_1",
            source="ocr",
            bbox=(0, 0, 90, 20),
            text="VECTOR EMBEDDINGS",
            confidence=0.9,
            font_size=18,
        )

        resolved_size = resolve_fallback_font_size_pt(block, scale_x=1.0, scale_y=1.0)

        self.assertLess(resolved_size, 18.0)
        self.assertGreaterEqual(resolved_size, 6.0)

    def test_resolve_fallback_font_size_pt_uses_tighter_limit_for_bracketed_label(self) -> None:
        block = TextBlock(
            id="ocr_fallback_font_2",
            source="ocr",
            bbox=(0, 0, 162.5, 38),
            text="[加權準確率]",
            confidence=0.9,
            font_size=29,
        )

        resolved_size = resolve_fallback_font_size_pt(block, scale_x=1.0, scale_y=1.0)

        self.assertLessEqual(resolved_size, 25.0)
        self.assertGreaterEqual(resolved_size, 6.0)

    def test_resolve_fallback_font_size_pt_uses_tighter_limit_for_tiny_latin_footer_label(self) -> None:
        block = TextBlock(
            id="ocr_fallback_font_3",
            source="ocr",
            bbox=(0, 0, 80.5, 13.5),
            text="NotebookLM",
            confidence=0.9,
            font_size=13,
        )

        resolved_size = resolve_fallback_font_size_pt(block, scale_x=1.0, scale_y=1.0)

        self.assertLessEqual(resolved_size, 11.0)
        self.assertGreaterEqual(resolved_size, 6.0)

    def test_resolve_fallback_font_size_pt_can_scale_up_title_like_single_line_ocr(self) -> None:
        block = TextBlock(
            id="ocr_fallback_font_4",
            source="ocr",
            bbox=(0, 0, 1184, 122.5),
            text="金融業生成式AI平台工程",
            confidence=0.9,
            font_size=48,
            block_role="body",
        )

        resolved_size = resolve_fallback_font_size_pt(block, scale_x=1.0, scale_y=1.0)

        self.assertGreater(resolved_size, 48.0)
        self.assertLessEqual(resolved_size, 96.0)

    def test_resolve_fallback_font_size_pt_uses_bold_measurement_for_latin_title(self) -> None:
        block = TextBlock(
            id="ocr_fallback_font_5",
            source="ocr",
            bbox=(438.0, 54.0, 609.5, 113.0),
            text="Agent",
            confidence=0.99,
            font_size=48.0,
            bold=True,
            block_role="title",
        )

        resolved_size = resolve_fallback_font_size_pt(block, scale_x=1.0, scale_y=1.0)

        self.assertLessEqual(resolved_size, 51.0)
        self.assertGreaterEqual(resolved_size, 48.0)

    @patch("pdf2ppt.ppt_render.fit_text_frame", return_value=True)
    @patch("pdf2ppt.ppt_render.resolve_fallback_font_size_pt", return_value=27.0)
    def test_add_text_block_clamps_oversized_fit_text_for_single_line_ocr(
        self,
        _resolve_fallback_font_size_pt: unittest.mock.Mock,
        _fit_text_frame: unittest.mock.Mock,
    ) -> None:
        class DummyFont:
            def __init__(self) -> None:
                self.name = None
                self.size = Pt(68)
                self.color = SimpleNamespace(rgb=None)
                self.bold = None
                self.italic = None

        class DummyRun:
            def __init__(self) -> None:
                self.text = ""
                self.font = DummyFont()

        class DummyParagraph:
            def __init__(self) -> None:
                self.alignment = None
                self._run = DummyRun()

            def add_run(self) -> DummyRun:
                return self._run

        class DummyTextFrame:
            def __init__(self) -> None:
                self.word_wrap = None
                self.vertical_anchor = None
                self.margin_left = None
                self.margin_right = None
                self.margin_top = None
                self.margin_bottom = None
                self.auto_size = None
                self.paragraphs = [DummyParagraph()]

            def clear(self) -> None:
                return None

        class DummyTextBox:
            def __init__(self) -> None:
                self.text_frame = DummyTextFrame()

        class DummyShapes:
            def __init__(self) -> None:
                self.last_textbox: DummyTextBox | None = None

            def add_textbox(self, *_args: object, **_kwargs: object) -> DummyTextBox:
                self.last_textbox = DummyTextBox()
                return self.last_textbox

        class DummySlide:
            def __init__(self) -> None:
                self.shapes = DummyShapes()

        slide = DummySlide()
        block = TextBlock(
            id="ocr_fit_4",
            source="ocr",
            bbox=(0, 0, 207.5, 71.5),
            text="97.5%",
            confidence=0.99,
            font_size=48.0,
            bold=True,
        )

        add_text_block(slide, block, scale_x=1.0, scale_y=1.0)

        textbox = slide.shapes.last_textbox
        self.assertIsNotNone(textbox)
        run = textbox.text_frame.paragraphs[0]._run
        self.assertEqual(run.font.size.pt, 27.0)

    def test_add_text_block_sets_east_asian_typeface_for_mixed_ocr_text(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        block = TextBlock(
            id="ocr_typeface_1",
            source="ocr",
            bbox=(0, 0, 102, 28.5),
            text="可營運AI",
            confidence=0.9,
            font_size=25.5,
        )

        add_text_block(slide, block, scale_x=1.0, scale_y=1.0)

        textbox = slide.shapes[-1]
        run_xml = textbox.text_frame.paragraphs[0].runs[0]._r.xml
        self.assertIn('a:latin typeface="Noto Sans CJK TC"', run_xml)
        self.assertIn('a:ea typeface="Noto Sans CJK TC"', run_xml)
        self.assertIn('a:cs typeface="Noto Sans CJK TC"', run_xml)
        paragraph_xml = textbox.text_frame.paragraphs[0]._p.xml
        self.assertIn('a:endParaRPr', paragraph_xml)
        self.assertIn('a:latin typeface="Noto Sans CJK TC"', paragraph_xml)
        self.assertIn('a:ea typeface="Noto Sans CJK TC"', paragraph_xml)
        self.assertIn('a:cs typeface="Noto Sans CJK TC"', paragraph_xml)

    def test_add_text_block_writes_wrap_none_for_single_line_ocr(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        block = TextBlock(
            id="ocr_wrap_xml_1",
            source="ocr",
            bbox=(0, 0, 167.5, 34.5),
            text="錯誤/不安全",
            confidence=0.9,
            font_size=29.5,
        )

        add_text_block(slide, block, scale_x=1.0, scale_y=1.0)

        textbox = slide.shapes[-1]
        self.assertIn('wrap="none"', textbox.text_frame._txBody.xml)

    def test_resolve_vertical_anchor_uses_middle_for_single_line_ocr(self) -> None:
        anchor = resolve_vertical_anchor(
            TextBlock(
                id="ocr_anchor_1",
                source="ocr",
                bbox=(0, 0, 100, 20),
                text="VECTOR EMBEDDINGS",
                confidence=0.9,
            )
        )
        self.assertEqual(int(anchor), 3)

    def test_resolve_vertical_anchor_uses_top_for_multiline_text(self) -> None:
        anchor = resolve_vertical_anchor(
            TextBlock(
                id="ocr_anchor_2",
                source="ocr",
                bbox=(0, 0, 100, 40),
                text="LINE 1\nLINE 2",
                confidence=0.9,
            )
        )
        self.assertEqual(int(anchor), 1)

    def test_resolve_ocr_fit_max_size_respects_single_line_width(self) -> None:
        size = resolve_ocr_fit_max_size(
            TextBlock(
                id="ocr_width_1",
                source="ocr",
                bbox=(0, 0, 120, 24),
                text="VECTOR EMBEDDINGS",
                confidence=0.9,
                font_size=12,
            ),
            font_path="/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            base_size=12,
            scale_x=1.0,
            scale_y=1.0,
            script="latin",
        )
        measured_width, _ = measure_text_dimensions(
            "VECTOR EMBEDDINGS",
            size,
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        )
        self.assertLessEqual(measured_width, 120 * 0.96)

    def test_resolve_component_telea_radius_uses_smaller_span_for_compact_wide_component(self) -> None:
        radius = inpainting_engines._resolve_component_telea_radius(
            width=169,
            height=57,
            base_radius=3.0,
            min_radius=1.8,
            max_radius=7.5,
            reference_span_px=48.0,
            edge_density=0.0317,
            edge_density_threshold=0.08,
            edge_density_min_factor=0.7,
        )

        self.assertLess(radius, 7.5)
        self.assertGreater(radius, 4.5)

    def test_clamp_isolated_label_telea_radius_caps_wide_low_edge_group(self) -> None:
        radius = inpainting_engines._clamp_isolated_label_telea_radius(
            7.5,
            width=217,
            height=70,
            edge_density=0.0155,
            group_size=1,
            proximity_px=0,
        )

        self.assertEqual(radius, 5.0)

    def test_clamp_isolated_label_telea_radius_keeps_other_groups_unchanged(self) -> None:
        radius = inpainting_engines._clamp_isolated_label_telea_radius(
            7.5,
            width=300,
            height=70,
            edge_density=0.0155,
            group_size=2,
            proximity_px=4,
        )

        self.assertEqual(radius, 7.5)

    def test_resolve_directional_inpaint_crop_bounds_prefers_vertical_context_for_compact_wide_component(self) -> None:
        gray = np.full((80, 120), 238, dtype=np.uint8)
        edges = np.zeros((80, 120), dtype=np.uint8)
        gray[20:50, 28:40] = 105
        gray[20:50, 80:92] = 118
        edges[20:50, 28:40] = 255
        edges[20:50, 80:92] = 255

        x0, y0, x1, y1 = inpainting_engines._resolve_directional_inpaint_crop_bounds(
            gray,
            edges,
            component_bbox=(40, 20, 40, 30),
            image_shape=gray.shape,
            base_padding=10,
        )

        self.assertEqual((y0, y1), (8, 62))
        self.assertEqual((x0, x1), (40, 80))

    def test_resolve_ocr_fit_max_size_uses_tighter_limit_for_bracketed_label(self) -> None:
        size = resolve_ocr_fit_max_size(
            TextBlock(
                id="ocr_width_2",
                source="ocr",
                bbox=(0, 0, 162.5, 38),
                text="[加權準確率]",
                confidence=0.9,
                font_size=29,
            ),
            font_path="/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            base_size=29,
            scale_x=1.0,
            scale_y=1.0,
            script="mixed",
        )

        self.assertLessEqual(size, 25)
        self.assertGreaterEqual(size, 6)

    def test_resolve_ocr_fit_max_size_uses_tighter_limit_for_tiny_latin_footer_label(self) -> None:
        size = resolve_ocr_fit_max_size(
            TextBlock(
                id="ocr_width_3",
                source="ocr",
                bbox=(0, 0, 80.5, 13.5),
                text="NotebookLM",
                confidence=0.9,
                font_size=13,
            ),
            font_path="/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            base_size=13,
            scale_x=1.0,
            scale_y=1.0,
            script="latin",
        )

        self.assertLessEqual(size, 11)
        self.assertGreaterEqual(size, 6)


class TextColorTests(unittest.TestCase):
    def test_estimate_text_color_prefers_dark_foreground(self) -> None:
        color = Image.new("RGB", (20, 12), (240, 242, 245))
        gray = color.convert("L")
        for x in range(4, 16):
            for y in range(3, 9):
                color.putpixel((x, y), (32, 92, 184))
                gray.putpixel((x, y), 70)
        self.assertEqual(estimate_text_color(color, gray), "#205CB8")

    def test_estimate_text_color_prefers_light_foreground(self) -> None:
        color = Image.new("RGB", (20, 12), (24, 28, 36))
        gray = color.convert("L")
        for x in range(4, 16):
            for y in range(3, 9):
                color.putpixel((x, y), (242, 244, 248))
                gray.putpixel((x, y), 240)
        self.assertEqual(estimate_text_color(color, gray), "#F2F4F8")


class TextStyleTests(unittest.TestCase):
    def test_estimate_text_bold_distinguishes_regular_and_bold(self) -> None:
        regular_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 34)
        bold_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 34)

        def render(font: ImageFont.FreeTypeFont) -> Image.Image:
            image = Image.new("L", (220, 70), 255)
            draw = ImageDraw.Draw(image)
            draw.text((10, 12), "VECTOR", font=font, fill=0)
            return image

        regular = render(regular_font)
        bold = render(bold_font)
        self.assertFalse(estimate_text_bold("VECTOR", regular))
        self.assertTrue(estimate_text_bold("VECTOR", bold))

    def test_promote_ocr_bold_blocks_marks_titles_bold(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_title_1",
                source="ocr",
                bbox=(0, 0, 300, 80),
                text="Main Title",
                confidence=0.9,
                font_size=42,
                block_role="title",
                bold=False,
            ),
            TextBlock(
                id="ocr_body_1",
                source="ocr",
                bbox=(0, 0, 200, 30),
                text="body text",
                confidence=0.9,
                font_size=18,
                block_role="body",
                bold=False,
            ),
        ]
        promote_ocr_bold_blocks(blocks)
        self.assertTrue(blocks[0].bold)
        self.assertFalse(blocks[1].bold)

    def test_promote_ocr_bold_blocks_does_not_promote_body_headers(self) -> None:
        blocks = [
            TextBlock(
                id="ocr_header_1",
                source="ocr",
                bbox=(0, 0, 220, 40),
                text="核心優勢",
                confidence=0.9,
                font_size=29.5,
                block_role="body",
                bold=False,
            ),
            TextBlock(
                id="ocr_body_2",
                source="ocr",
                bbox=(0, 0, 300, 35),
                text="這是一段較長的正文內容",
                confidence=0.9,
                font_size=27.0,
                block_role="body",
                bold=False,
            ),
            TextBlock(
                id="ocr_body_3",
                source="ocr",
                bbox=(0, 0, 300, 35),
                text="普通內文",
                confidence=0.9,
                font_size=18.0,
                block_role="body",
                bold=False,
            ),
        ]
        promote_ocr_bold_blocks(blocks)
        self.assertFalse(blocks[0].bold)
        self.assertFalse(blocks[1].bold)
        self.assertFalse(blocks[2].bold)


if __name__ == "__main__":
    unittest.main()
