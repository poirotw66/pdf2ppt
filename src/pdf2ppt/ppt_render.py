from __future__ import annotations

import io
import logging
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from .models import TextBlock
from .text_style import (
    choose_measurement_font,
    classify_text_script,
    default_font_family,
    default_ppt_typefaces,
    measure_text_dimensions,
    ocr_fit_height_cap_ratio,
    single_line_fit_width_ratio,
)

logger = logging.getLogger(__name__)

EMU_PER_PT = 12700


def render_page_to_slide(
    presentation: Any,
    blank_layout: Any,
    page_result: Any,
    *,
    slide_width_pt: float,
    slide_height_pt: float,
) -> None:
    slide = presentation.slides.add_slide(blank_layout)
    scale_x = slide_width_pt / page_result.width_pt
    scale_y = slide_height_pt / page_result.height_pt

    if page_result.background_image_bytes is not None:
        slide.shapes.add_picture(
            io.BytesIO(page_result.background_image_bytes),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    for image_element in page_result.image_elements:
        left, top, width, height = bbox_to_shape_geometry(image_element.bbox, scale_x, scale_y)
        slide.shapes.add_picture(
            io.BytesIO(image_element.png_bytes),
            left,
            top,
            width=width,
            height=height,
        )

    for block in sorted(page_result.text_blocks, key=lambda item: item.reading_order):
        add_text_block(slide, block, scale_x=scale_x, scale_y=scale_y)


def add_text_block(slide: Any, block: TextBlock, *, scale_x: float, scale_y: float) -> None:
    script = classify_text_script(block.text)
    resolved_font_family = block.font_family or default_font_family(script)
    resolved_typefaces = default_ppt_typefaces(script)
    if block.font_family:
        resolved_typefaces = {"latin": block.font_family, "ea": block.font_family, "cs": block.font_family}
    left, top, width, height = bbox_to_shape_geometry(block.bbox, scale_x, scale_y)
    textbox = slide.shapes.add_textbox(left, top, max(width, Emu(1)), max(height, Emu(1)))
    text_frame = textbox.text_frame
    text_frame.word_wrap = should_wrap_text_block(block)
    text_frame.vertical_anchor = resolve_vertical_anchor(block)
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.clear()

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = block.text
    used_fit_text = fit_text_frame(text_frame, block, scale_x=scale_x, scale_y=scale_y)

    font = run.font
    font.name = resolved_font_family
    _set_run_typefaces(run, resolved_typefaces)
    if used_fit_text and block.source == "ocr" and "\n" not in block.text:
        safe_font_size_pt = resolve_fallback_font_size_pt(block, scale_x=scale_x, scale_y=scale_y)
        current_font_size_pt = getattr(font.size, "pt", None) if font.size is not None else None
        if current_font_size_pt is None or current_font_size_pt > safe_font_size_pt:
            font.size = Pt(safe_font_size_pt)
    if block.font_size and not used_fit_text:
        font.size = Pt(resolve_fallback_font_size_pt(block, scale_x=scale_x, scale_y=scale_y))
    if block.font_color:
        font.color.rgb = RGBColor.from_string(block.font_color.lstrip("#"))
    font.bold = block.bold
    font.italic = block.italic
    _set_paragraph_end_typefaces(paragraph, resolved_typefaces)


def fit_text_frame(text_frame: Any, block: TextBlock, *, scale_x: float, scale_y: float) -> bool:
    if block.source != "ocr" or not block.text.strip():
        return False
    script = classify_text_script(block.text)
    font_path = choose_measurement_font(script)
    if font_path is None:
        logger.debug("No measurement font available for script=%s; skipping fit_text", script)
        return False

    font_family = block.font_family or default_font_family(script)
    base_size = max(6, int(round((block.font_size or 12.0) * min(scale_x, scale_y))))
    max_size = max(
        6,
        resolve_ocr_fit_max_size(
            block,
            font_path=font_path,
            base_size=base_size,
            scale_x=scale_x,
            scale_y=scale_y,
            script=script,
        ),
    )
    try:
        text_frame.fit_text(
            font_family=font_family,
            max_size=max_size,
            bold=block.bold,
            italic=block.italic,
            font_file=font_path,
        )
        return True
    except TypeError:
        logger.debug("python-pptx fit_text rejected OCR block %s", block.id)
        return False


def resolve_fallback_font_size_pt(block: TextBlock, *, scale_x: float, scale_y: float) -> float:
    base_size_pt = max(6.0, (block.font_size or 12.0) * min(scale_x, scale_y))
    if not block.text.strip():
        return base_size_pt

    script = classify_text_script(block.text)
    font_path = choose_measurement_font(script)
    if font_path is None:
        return base_size_pt

    width_pt = max(1.0, (block.bbox[2] - block.bbox[0]) * scale_x)
    height_pt = max(1.0, (block.bbox[3] - block.bbox[1]) * scale_y)
    width_limit = width_pt if "\n" in block.text else resolve_single_line_width_limit_pt(
        block.text,
        width_pt,
        script,
        height_pt=height_pt,
    )
    start_size = max(6, int(round(base_size_pt)))
    if is_title_like_single_line_ocr_block(block, width_pt=width_pt, height_pt=height_pt):
        start_size = max(
            start_size,
            resolve_ocr_fit_max_size(
                block,
                font_path=font_path,
                base_size=start_size,
                scale_x=scale_x,
                scale_y=scale_y,
                script=script,
            ),
        )
    for size in range(start_size, 5, -1):
        measured_width, measured_height = measure_text_dimensions(block.text, size, font_path)
        if measured_width <= width_limit and measured_height <= height_pt * 1.03:
            return float(size)
    return 6.0


def should_wrap_text_block(block: TextBlock) -> bool:
    if block.source == "ocr" and "\n" not in block.text:
        return False
    return True


def resolve_vertical_anchor(block: TextBlock) -> MSO_VERTICAL_ANCHOR:
    if block.source == "ocr" and "\n" not in block.text:
        return MSO_VERTICAL_ANCHOR.MIDDLE
    return MSO_VERTICAL_ANCHOR.TOP


def resolve_ocr_fit_max_size(
    block: TextBlock,
    *,
    font_path: str,
    base_size: int,
    scale_x: float,
    scale_y: float,
    script: str,
) -> int:
    height_pt = max(1.0, (block.bbox[3] - block.bbox[1]) * scale_y)
    max_size = min(96, max(base_size + 3, int(round(height_pt * ocr_fit_height_cap_ratio(script)))))
    if "\n" in block.text:
        return max_size

    width_pt = max(1.0, (block.bbox[2] - block.bbox[0]) * scale_x)
    height_pt = max(1.0, (block.bbox[3] - block.bbox[1]) * scale_y)
    width_limit = resolve_single_line_width_limit_pt(block.text, width_pt, script, height_pt=height_pt)
    size_cap = max(6, max_size)
    for size in range(size_cap, 5, -1):
        measured_width, measured_height = measure_text_dimensions(block.text, size, font_path)
        if measured_width <= width_limit and measured_height <= height_pt * 1.03:
            return size
    return min(base_size, max_size)


def resolve_single_line_width_limit_pt(text: str, width_pt: float, script: str, *, height_pt: float) -> float:
    ratio = single_line_fit_width_ratio(script)
    if is_bracketed_ocr_label(text):
        ratio = min(ratio, 0.87)
    if is_tiny_latin_footer_label(text, script, width_pt, height_pt):
        ratio = min(ratio, 0.90)
    return width_pt * ratio


def is_bracketed_ocr_label(text: str) -> bool:
    stripped = text.strip()
    if len(stripped) < 4 or len(stripped) > 12:
        return False
    bracket_pairs = {
        "[": "]",
        "(": ")",
        "{": "}",
        "【": "】",
        "（": "）",
        "［": "］",
    }
    closing = bracket_pairs.get(stripped[0])
    return closing is not None and stripped.endswith(closing)


def is_tiny_latin_footer_label(text: str, script: str, width_pt: float, height_pt: float) -> bool:
    stripped = text.strip()
    if script != "latin":
        return False
    if len(stripped) < 6 or len(stripped) > 16:
        return False
    return width_pt <= 100.0 and height_pt <= 16.0


def is_title_like_single_line_ocr_block(block: TextBlock, *, width_pt: float, height_pt: float) -> bool:
    if block.source != "ocr" or "\n" in block.text:
        return False
    if block.block_role == "title":
        return True
    stripped = block.text.strip()
    if len(stripped) < 4 or len(stripped) > 36:
        return False
    if (block.font_size or 0.0) < 28.0:
        return False
    return width_pt >= 400.0 and height_pt >= 56.0


def bbox_to_shape_geometry(
    bbox: tuple[float, float, float, float],
    scale_x: float,
    scale_y: float,
) -> tuple[Emu, Emu, Emu, Emu]:
    x0, y0, x1, y1 = bbox
    left = pt_to_emu(x0 * scale_x)
    top = pt_to_emu(y0 * scale_y)
    width = pt_to_emu(max(1.0, (x1 - x0) * scale_x))
    height = pt_to_emu(max(1.0, (y1 - y0) * scale_y))
    return left, top, width, height


def _set_run_typefaces(run: Any, typefaces: dict[str, str]) -> None:
    if not hasattr(run, "_r"):
        return
    r_pr = run._r.get_or_add_rPr()
    for child_tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
        for child in list(r_pr):
            if child.tag == child_tag:
                r_pr.remove(child)
    for tag_name, key in (("a:latin", "latin"), ("a:ea", "ea"), ("a:cs", "cs")):
        typeface = OxmlElement(tag_name)
        typeface.set("typeface", typefaces[key])
        r_pr.append(typeface)


def _set_paragraph_end_typefaces(paragraph: Any, typefaces: dict[str, str]) -> None:
    if not hasattr(paragraph, "_p"):
        return
    end_paragraph_properties = paragraph._p.get_or_add_endParaRPr()
    for child_tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
        for child in list(end_paragraph_properties):
            if child.tag == child_tag:
                end_paragraph_properties.remove(child)
    for tag_name, key in (("a:latin", "latin"), ("a:ea", "ea"), ("a:cs", "cs")):
        typeface = OxmlElement(tag_name)
        typeface.set("typeface", typefaces[key])
        end_paragraph_properties.append(typeface)


def pt_to_emu(value: float) -> Emu:
    return Emu(int(round(value * EMU_PER_PT)))
