from __future__ import annotations

import json
import logging
from time import perf_counter
from typing import Any, Callable

import fitz
from PIL import Image
from pptx import Presentation

from .background import (
    BackgroundInpaintingEngine,
    BackgroundInpaintingError,
    BackgroundRenderResult,
    OpenCvFastInpaintingEngine,
    WhiteBoxInpaintingEngine,
    build_mask_shapes,
    build_text_mask_image,
    compute_mask_crop_box,
    estimate_background_complexity,
    mask_area_ratio,
    mask_text_regions_with_white_boxes,
    render_overlay_background,
    resolve_background_inpainting_engine,
    write_debug_artifacts,
)
from .block_analysis import (
    bbox_area,
    choose_background_mode,
    classify_page,
    compute_page_signals,
    enrich_ocr_blocks,
    filter_suspicious_ocr_blocks,
    intersection_ratio,
    map_blocks_to_page_coordinates,
    resolve_background_render_dpi,
    resolve_render_dpi,
    safe_crop,
    score_page,
    select_text_blocks,
)
from .core import (
    ConversionOptions,
    OcrInitializationError,
    OcrPageData,
    OcrProcessingError,
    PageSignals,
)
from .models import ConversionReport, ImagePlacement, PageResult, QualityScore, TextBlock
from .native_extraction import (
    assign_block_roles,
    extract_image_elements,
    extract_native_text_blocks,
    most_common_or_none,
    pil_to_image_bytes,
    pil_to_png_bytes,
    promote_ocr_bold_blocks,
    sort_text_blocks,
)
from .ocr import OcrEngine
from .page_analysis import render_page_image
from .ppt_render import (
    add_text_block,
    bbox_to_shape_geometry,
    fit_text_frame,
    pt_to_emu,
    render_page_to_slide,
    resolve_vertical_anchor,
    resolve_ocr_fit_max_size,
    should_wrap_text_block,
)
from .text_style import (
    build_text_fit_debug_entry,
    choose_measurement_font,
    classify_text_script,
    default_font_family,
    estimate_font_size,
    estimate_text_bold,
    estimate_text_color,
    extract_text_foreground_mask,
    is_cjk,
    measure_text_dimensions,
    script_target_height_ratio,
    script_target_width_ratio,
    single_line_fit_width_ratio,
    write_text_fit_debug_report,
)

logger = logging.getLogger(__name__)

ProgressCallback = Callable[[int, int], None]


def convert_pdf(
    options: ConversionOptions,
    *,
    progress_callback: ProgressCallback | None = None,
) -> ConversionReport:
    conversion_started_at = perf_counter()
    options.output_path.parent.mkdir(parents=True, exist_ok=True)
    options.report_path.parent.mkdir(parents=True, exist_ok=True)

    logger.info("Starting PDF conversion: %s -> %s", options.input_path, options.output_path)
    ocr_engine = OcrEngine(
        lang=options.lang,
        model_root=options.ocr_model_root,
        use_doc_orientation=options.ocr_use_doc_orientation,
        use_textline_orientation=options.ocr_use_textline_orientation,
        use_doc_unwarping=options.use_doc_unwarping,
        det_thresh=options.ocr_det_thresh,
        det_box_thresh=options.ocr_det_box_thresh,
        drop_score=options.ocr_drop_score,
    )
    page_results: list[PageResult] = []

    with fitz.open(options.input_path) as document:
        presentation = Presentation()
        first_page = document[0]
        presentation.slide_width = pt_to_emu(first_page.rect.width)
        presentation.slide_height = pt_to_emu(first_page.rect.height)
        blank_layout = presentation.slide_layouts[6]
        total_pages = document.page_count
        if progress_callback is not None:
            progress_callback(0, total_pages)

        for page_index in range(total_pages):
            page = document[page_index]
            logger.info("Processing page %s/%s", page_index + 1, total_pages)
            try:
                page_result = analyze_page(page, options, ocr_engine)
            except Exception as error:
                logger.exception("Failed to process page %s", page.number + 1)
                raise RuntimeError(f"Failed to process page {page.number + 1}: {error}") from error
            page_results.append(page_result)
            slide_render_started_at = perf_counter()
            render_page_to_slide(
                presentation,
                blank_layout,
                page_result,
                slide_width_pt=first_page.rect.width,
                slide_height_pt=first_page.rect.height,
            )
            slide_render_seconds = perf_counter() - slide_render_started_at
            logger.info("Page %s slide render: %.3fs", page.number + 1, slide_render_seconds)
            if progress_callback is not None:
                progress_callback(page_index + 1, total_pages)

        presentation_save_started_at = perf_counter()
        presentation.save(options.output_path)
        presentation_save_seconds = perf_counter() - presentation_save_started_at
        logger.info("Presentation save: %.3fs", presentation_save_seconds)

    report = ConversionReport(
        input_path=str(options.input_path),
        output_path=str(options.output_path),
        pages=page_results,
    )
    options.report_path.write_text(
        json.dumps(report.to_dict(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    logger.info(
        "Finished PDF conversion with %s page(s) in %.3fs",
        len(page_results),
        perf_counter() - conversion_started_at,
    )
    return report


def analyze_page(page: fitz.Page, options: ConversionOptions, ocr_engine: OcrEngine) -> PageResult:
    page_started_at = perf_counter()
    native_extract_seconds = 0.0
    signal_seconds = 0.0
    ocr_render_seconds = 0.0
    ocr_seconds = 0.0
    ocr_enrich_seconds = 0.0
    page_mode_seconds = 0.0
    background_render_seconds = 0.0
    background_process_seconds = 0.0
    background_encode_seconds = 0.0

    native_extract_started_at = perf_counter()
    native_blocks, image_boxes = extract_native_text_blocks(page)
    native_extract_seconds = perf_counter() - native_extract_started_at
    signal_started_at = perf_counter()
    signals = compute_page_signals(page, native_blocks, image_boxes)
    signal_seconds = perf_counter() - signal_started_at
    page_kind = classify_page(signals)
    render_dpi = resolve_render_dpi(options)
    background_render_dpi = resolve_background_render_dpi(options)

    page_number = page.number + 1
    has_approved_ocr_blocks = page_number in options.approved_ocr_blocks_by_page
    need_ocr = page_kind in {"scanned", "hybrid"} or not native_blocks
    ocr_blocks: list[TextBlock] = []
    page_image: Any | None = None
    ocr_reference_image: Any | None = None
    if has_approved_ocr_blocks:
        ocr_render_started_at = perf_counter()
        page_image = render_page_image(page, dpi=render_dpi)
        ocr_render_seconds = perf_counter() - ocr_render_started_at
        ocr_reference_image = page_image
        ocr_enrich_started_at = perf_counter()
        ocr_blocks = resolve_approved_ocr_blocks(
            page_number=page_number,
            options=options,
            page_rect=page.rect,
            page_image=page_image,
            ocr_engine=ocr_engine,
        )
        ocr_enrich_seconds = perf_counter() - ocr_enrich_started_at
    elif need_ocr:
        ocr_render_started_at = perf_counter()
        page_image = render_page_image(page, dpi=render_dpi)
        ocr_render_seconds = perf_counter() - ocr_render_started_at

        ocr_started_at = perf_counter()
        ocr_page_data = ocr_engine.extract_text_blocks(page_image, page.number + 1)
        ocr_seconds = perf_counter() - ocr_started_at
        ocr_enrich_started_at = perf_counter()
        ocr_reference_image = ocr_page_data.image
        ocr_blocks = enrich_ocr_blocks(ocr_page_data.blocks, ocr_reference_image)
        ocr_blocks = map_blocks_to_page_coordinates(ocr_blocks, ocr_reference_image.size, page.rect)
        filtered_ocr_blocks = filter_suspicious_ocr_blocks(ocr_blocks, page.rect)
        removed_block_count = len(ocr_blocks) - len(filtered_ocr_blocks)
        if removed_block_count > 0:
            logger.info(
                "Filtered %s suspicious OCR block(s) on page %s",
                removed_block_count,
                page.number + 1,
            )
        ocr_blocks = filtered_ocr_blocks
        ocr_enrich_seconds = perf_counter() - ocr_enrich_started_at

    page_mode_started_at = perf_counter()
    text_blocks = select_text_blocks(page_kind, native_blocks, ocr_blocks)
    quality = score_page(page_kind, text_blocks, native_blocks, ocr_blocks)
    background_mode, fallback_reason = choose_background_mode(
        page_kind=page_kind,
        quality=quality,
        has_text=bool(text_blocks),
        has_visuals=bool(image_boxes) or signals.drawing_count > 0,
    )
    page_mode_seconds = perf_counter() - page_mode_started_at
    logger.debug(
        "Page %s classified as %s with background mode %s",
        page.number + 1,
        page_kind,
        background_mode,
    )

    background_image_bytes: bytes | None = None
    image_elements: list[ImagePlacement] = []
    background_inpaint_engine: str | None = None
    background_inpaint_note: str | None = None
    mask_image: Any | None = None
    if background_mode == "elements":
        image_elements = extract_image_elements(page, image_boxes, options.render_dpi)
    else:
        has_ocr_reference_image = ocr_reference_image is not None
        should_downscale_background = background_render_dpi < render_dpi
        uses_original_page_geometry = not options.ocr_use_doc_orientation and not options.use_doc_unwarping
        can_reuse_ocr_raster_for_background = (
            has_ocr_reference_image and should_downscale_background and uses_original_page_geometry
        )
        if page_image is None and background_render_dpi == render_dpi:
            background_render_started_at = perf_counter()
            page_image = render_page_image(page, dpi=render_dpi)
            background_render_seconds += perf_counter() - background_render_started_at
        if ocr_reference_image is not None and background_render_dpi == render_dpi:
            background_image = ocr_reference_image
        elif can_reuse_ocr_raster_for_background:
            background_render_started_at = perf_counter()
            background_image = resize_rendered_page_image(
                ocr_reference_image,
                source_dpi=render_dpi,
                target_dpi=background_render_dpi,
            )
            background_render_seconds += perf_counter() - background_render_started_at
        else:
            background_render_started_at = perf_counter()
            background_image = render_page_image(page, dpi=background_render_dpi)
            background_render_seconds += perf_counter() - background_render_started_at
        if background_mode == "overlay" and text_blocks:
            mask_blocks = [block for block in text_blocks if block.source == "ocr"] or text_blocks
            background_process_started_at = perf_counter()
            background_result = render_overlay_background(
                background_image,
                mask_blocks,
                page.rect,
                options=options,
            )
            background_process_seconds = perf_counter() - background_process_started_at
            background_image = background_result.image
            background_inpaint_engine = background_result.engine_name
            background_inpaint_note = background_result.note
            mask_image = background_result.mask_image
        background_encode_started_at = perf_counter()
        background_image_bytes = pil_to_image_bytes(
            background_image,
            image_format=options.background_image_format,
            jpeg_quality=options.background_jpeg_quality,
        )
        background_encode_seconds = perf_counter() - background_encode_started_at
    if options.debug_dir is not None and ocr_blocks:
        debug_blocks = [block for block in text_blocks if block.source == "ocr"] or ocr_blocks
        debug_masked_image = background_image if background_mode == "overlay" else ocr_reference_image
        if debug_masked_image is not None and ocr_reference_image is not None:
            if debug_masked_image.size != ocr_reference_image.size:
                debug_masked_image = ocr_reference_image
        write_debug_artifacts(
            debug_dir=options.debug_dir,
            page_number=page_number,
            page_image=ocr_reference_image,
            masked_image=debug_masked_image,
            mask_image=mask_image,
            text_blocks=debug_blocks,
            page_rect=page.rect,
            engine_name=background_inpaint_engine,
            engine_note=background_inpaint_note,
            extra_images=background_result.debug_images if background_mode == "overlay" else None,
        )
        write_text_fit_debug_report(
            debug_dir=options.debug_dir,
            page_number=page_number,
            text_blocks=debug_blocks,
        )

    page_total_seconds = perf_counter() - page_started_at
    logger.info(
        "Page %s timings: native_extract=%.3fs signals=%.3fs ocr_render=%.3fs ocr=%.3fs ocr_enrich=%.3fs page_mode=%.3fs background_render=%.3fs background_process=%.3fs background_encode=%.3fs total=%.3fs",
        page_number,
        native_extract_seconds,
        signal_seconds,
        ocr_render_seconds,
        ocr_seconds,
        ocr_enrich_seconds,
        page_mode_seconds,
        background_render_seconds,
        background_process_seconds,
        background_encode_seconds,
        page_total_seconds,
    )

    return PageResult(
        page_number=page_number,
        page_kind=page_kind,
        background_mode=background_mode,
        width_pt=page.rect.width,
        height_pt=page.rect.height,
        text_blocks=text_blocks if background_mode != "full-page" else [],
        quality_score=quality,
        fallback_reason=fallback_reason,
        background_inpaint_engine=background_inpaint_engine,
        background_inpaint_note=background_inpaint_note,
        background_image_bytes=background_image_bytes,
        image_elements=image_elements,
    )


def resolve_approved_ocr_blocks(
    *,
    page_number: int,
    options: ConversionOptions,
    page_rect: fitz.Rect,
    page_image: Image.Image,
    ocr_engine: OcrEngine,
) -> list[TextBlock]:
    approved_blocks = options.approved_ocr_blocks_by_page.get(page_number, [])
    if not approved_blocks:
        return []

    source_image_size = options.approved_ocr_image_size_by_page.get(page_number, page_image.size)
    scaled_blocks = scale_blocks_to_image_size(approved_blocks, source_image_size, page_image.size)
    resolved_blocks = recognize_missing_approved_block_texts(scaled_blocks, page_image=page_image, ocr_engine=ocr_engine, page_number=page_number)
    resolved_blocks = [block for block in resolved_blocks if block.text.strip()]
    if not resolved_blocks:
        return []
    enriched_blocks = enrich_ocr_blocks(resolved_blocks, page_image)
    return map_blocks_to_page_coordinates(enriched_blocks, page_image.size, page_rect)


def recognize_missing_approved_block_texts(
    blocks: list[TextBlock],
    *,
    page_image: Image.Image,
    ocr_engine: OcrEngine,
    page_number: int,
) -> list[TextBlock]:
    resolved_blocks: list[TextBlock] = []
    for index, block in enumerate(blocks, start=1):
        if block.text.strip():
            resolved_blocks.append(block)
            continue
        recognized_text, recognized_confidence = ocr_engine.recognize_text_in_box(
            page_image,
            block.image_bbox or block.bbox,
            page_number=page_number * 1000 + index,
        )
        resolved_blocks.append(
            TextBlock(
                id=block.id,
                source=block.source,
                bbox=block.bbox,
                text=recognized_text,
                confidence=recognized_confidence if recognized_text else block.confidence,
                font_family=block.font_family,
                font_size=block.font_size,
                font_color=block.font_color,
                bold=block.bold,
                italic=block.italic,
                reading_order=block.reading_order,
                block_role=block.block_role,
                image_bbox=block.image_bbox,
                image_polygon=block.image_polygon,
            )
        )
    return resolved_blocks


def scale_blocks_to_image_size(
    blocks: list[TextBlock],
    source_size: tuple[int, int],
    target_size: tuple[int, int],
) -> list[TextBlock]:
    source_width, source_height = source_size
    target_width, target_height = target_size
    if source_width <= 0 or source_height <= 0:
        return blocks
    scale_x = target_width / source_width
    scale_y = target_height / source_height
    if abs(scale_x - 1.0) < 1e-6 and abs(scale_y - 1.0) < 1e-6:
        return blocks

    scaled_blocks: list[TextBlock] = []
    for block in blocks:
        x0, y0, x1, y1 = block.bbox
        scaled_blocks.append(
            TextBlock(
                id=block.id,
                source=block.source,
                bbox=(x0 * scale_x, y0 * scale_y, x1 * scale_x, y1 * scale_y),
                text=block.text,
                confidence=block.confidence,
                font_family=block.font_family,
                font_size=block.font_size * scale_y if block.font_size else None,
                font_color=block.font_color,
                bold=block.bold,
                italic=block.italic,
                reading_order=block.reading_order,
                block_role=block.block_role,
                image_bbox=(
                    (
                        block.image_bbox[0] * scale_x,
                        block.image_bbox[1] * scale_y,
                        block.image_bbox[2] * scale_x,
                        block.image_bbox[3] * scale_y,
                    )
                    if block.image_bbox is not None
                    else None
                ),
                image_polygon=(
                    tuple((point[0] * scale_x, point[1] * scale_y) for point in block.image_polygon)
                    if block.image_polygon is not None
                    else None
                ),
            )
        )
    return scaled_blocks


def resize_rendered_page_image(image: Image.Image, *, source_dpi: int, target_dpi: int) -> Image.Image:
    if target_dpi >= source_dpi:
        return image

    scale = target_dpi / max(source_dpi, 1)
    target_width = max(1, int(round(image.width * scale)))
    target_height = max(1, int(round(image.height * scale)))
    return image.resize((target_width, target_height), resample=Image.Resampling.LANCZOS)


__all__ = [
    "BackgroundInpaintingEngine",
    "BackgroundInpaintingError",
    "BackgroundRenderResult",
    "ConversionOptions",
    "ConversionReport",
    "ImagePlacement",
    "OcrEngine",
    "OcrInitializationError",
    "OcrPageData",
    "OcrProcessingError",
    "OpenCvFastInpaintingEngine",
    "PageResult",
    "PageSignals",
    "ProgressCallback",
    "QualityScore",
    "TextBlock",
    "WhiteBoxInpaintingEngine",
    "add_text_block",
    "analyze_page",
    "assign_block_roles",
    "bbox_area",
    "bbox_to_shape_geometry",
    "build_mask_shapes",
    "build_text_fit_debug_entry",
    "build_text_mask_image",
    "choose_background_mode",
    "choose_measurement_font",
    "classify_page",
    "classify_text_script",
    "compute_mask_crop_box",
    "compute_page_signals",
    "convert_pdf",
    "default_font_family",
    "enrich_ocr_blocks",
    "estimate_background_complexity",
    "estimate_font_size",
    "estimate_text_bold",
    "estimate_text_color",
    "extract_image_elements",
    "extract_native_text_blocks",
    "extract_text_foreground_mask",
    "filter_suspicious_ocr_blocks",
    "fit_text_frame",
    "intersection_ratio",
    "is_cjk",
    "map_blocks_to_page_coordinates",
    "mask_area_ratio",
    "mask_text_regions_with_white_boxes",
    "measure_text_dimensions",
    "most_common_or_none",
    "pil_to_png_bytes",
    "promote_ocr_bold_blocks",
    "pil_to_image_bytes",
    "pt_to_emu",
    "render_overlay_background",
    "render_page_image",
    "render_page_to_slide",
    "resolve_background_inpainting_engine",
    "resolve_background_render_dpi",
    "resolve_ocr_fit_max_size",
    "resolve_vertical_anchor",
    "resolve_render_dpi",
    "safe_crop",
    "score_page",
    "script_target_height_ratio",
    "script_target_width_ratio",
    "select_text_blocks",
    "should_wrap_text_block",
    "single_line_fit_width_ratio",
    "sort_text_blocks",
    "write_debug_artifacts",
    "write_text_fit_debug_report",
]
